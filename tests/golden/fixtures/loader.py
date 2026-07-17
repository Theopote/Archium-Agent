"""Layer 2: real fixture acceptance case loader."""

from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from archium.application.ingestion_service import IngestionService
from archium.application.presentation_models import PresentationRequest
from archium.domain.enums import ProjectType
from archium.domain.fact import ProjectFact
from archium.domain.project import Project
from archium.infrastructure.database.repositories import FactRepository, ProjectRepository
from sqlalchemy.orm import Session

_MANIFESTS_DIR = Path(__file__).resolve().parent / "manifests"
_FILES_DIR = Path(__file__).resolve().parent / "files"


@dataclass(frozen=True)
class FixtureCase:
    id: str
    name: str
    project_name: str
    project_type: ProjectType
    request: PresentationRequest
    expectations: dict[str, Any]
    export_presentation_spec: bool
    source_files: tuple[Path, ...]


def list_fixture_manifest_paths() -> list[Path]:
    return sorted(_MANIFESTS_DIR.glob("*.fixture.json"))


def load_fixture_case(path: Path) -> FixtureCase:
    payload = json.loads(path.read_text(encoding="utf-8"))
    project = payload["project"]
    request_data = payload["request"]
    from archium.domain.enums import PresentationType

    request = PresentationRequest(
        title=request_data["title"],
        audience=request_data["audience"],
        purpose=request_data["purpose"],
        duration_minutes=int(request_data.get("duration_minutes", 20)),
        target_slide_count=int(request_data.get("target_slide_count", 4)),
        core_message=request_data.get("core_message", ""),
        required_sections=list(request_data.get("required_sections", [])),
        presentation_type=PresentationType(request_data.get("presentation_type", "client_review")),
    )
    expectations = dict(payload.get("expectations", {}))
    export_spec = bool(expectations.pop("export_presentation_spec", False))

    source_files: list[Path] = []
    for entry in payload.get("files", []):
        relative = Path(str(entry["relative_path"]))
        resolved = (_FILES_DIR / relative).resolve()
        required = bool(entry.get("required", False))
        if resolved.is_file():
            source_files.append(resolved)
        elif required:
            raise FileNotFoundError(f"Required fixture file missing: {resolved}")

    return FixtureCase(
        id=str(payload["id"]),
        name=str(payload["name"]),
        project_name=str(project["name"]),
        project_type=ProjectType(str(project["project_type"])),
        request=request,
        expectations=expectations,
        export_presentation_spec=export_spec,
        source_files=tuple(source_files),
    )


def materialize_inline_docx(path: Path, payload: dict[str, Any]) -> Path:
    """Create a DOCX from manifest inline paragraphs when bundled files are absent."""
    from docx import Document

    path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    for paragraph in payload.get("paragraphs", []):
        document.add_paragraph(str(paragraph))
    document.save(path)
    return path


def materialize_inline_xlsx(path: Path, payload: dict[str, Any]) -> Path:
    """Create an XLSX from manifest inline rows when bundled files are absent."""
    from openpyxl import Workbook

    path.parent.mkdir(parents=True, exist_ok=True)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = str(payload.get("sheet_name", "数据"))
    for row_index, row in enumerate(payload.get("rows", []), start=1):
        for col_index, value in enumerate(row, start=1):
            sheet.cell(row=row_index, column=col_index, value=value)
    workbook.save(path)
    return path


def materialize_inline_pdf(path: Path, payload: dict[str, Any]) -> Path:
    """Create a minimal PDF with real parser coverage when bundled files are absent."""
    import fitz

    path.parent.mkdir(parents=True, exist_ok=True)
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), str(payload.get("text", "测试 PDF 内容")))
    document.save(path)
    document.close()
    return path


def materialize_inline_pptx(path: Path, payload: dict[str, Any]) -> Path:
    """Create a minimal PPTX when bundled files are absent."""
    from pptx import Presentation

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = str(payload.get("title", "测试汇报"))
    body = slide.placeholders[1]
    body.text = str(payload.get("body", "由 fixture 生成的测试内容"))
    presentation.save(path)
    return path


def materialize_inline_image(path: Path, payload: dict[str, Any]) -> Path:
    """Create a low-resolution image for parser smoke coverage."""
    from PIL import Image

    path.parent.mkdir(parents=True, exist_ok=True)
    width = int(payload.get("width", 48))
    height = int(payload.get("height", 36))
    image = Image.new("RGB", (width, height), color=tuple(payload.get("color", [180, 170, 160])))
    image.save(path)
    return path


def resolve_fixture_scratch_dir(payload: dict[str, Any], scratch_dir: Path) -> Path:
    """Support Chinese and spaced path prefixes for cross-platform parser tests."""
    prefix = str(payload.get("path_prefix", "")).strip()
    if not prefix:
        return scratch_dir
    return scratch_dir / Path(prefix)


def materialize_inline_fallbacks(
    payload: dict[str, Any],
    scratch_dir: Path,
) -> list[Path]:
    """Build parser-ready files from inline_* blocks when no real files exist."""
    base_dir = resolve_fixture_scratch_dir(payload, scratch_dir)
    paths: list[Path] = []
    inline_docx = payload.get("inline_docx")
    if inline_docx is not None:
        paths.append(
            materialize_inline_docx(
                base_dir / str(inline_docx.get("filename", "inline.docx")),
                inline_docx,
            )
        )
    inline_xlsx = payload.get("inline_xlsx")
    if inline_xlsx is not None:
        paths.append(
            materialize_inline_xlsx(
                base_dir / str(inline_xlsx.get("filename", "inline.xlsx")),
                inline_xlsx,
            )
        )
    inline_pdf = payload.get("inline_pdf")
    if inline_pdf is not None:
        paths.append(
            materialize_inline_pdf(
                base_dir / str(inline_pdf.get("filename", "inline.pdf")),
                inline_pdf,
            )
        )
    inline_pptx = payload.get("inline_pptx")
    if inline_pptx is not None:
        paths.append(
            materialize_inline_pptx(
                base_dir / str(inline_pptx.get("filename", "inline.pptx")),
                inline_pptx,
            )
        )
    inline_image = payload.get("inline_image")
    if inline_image is not None:
        paths.append(
            materialize_inline_image(
                base_dir / str(inline_image.get("filename", "inline.png")),
                inline_image,
            )
        )
    return paths


def seed_fixture_case(
    session: Session,
    path: Path,
    *,
    scratch_dir: Path,
) -> tuple[FixtureCase, Project, list[Path]]:
    """Import real files through IngestionService (real parsers)."""
    payload = json.loads(path.read_text(encoding="utf-8"))
    case = load_fixture_case(path)
    project = ProjectRepository(session).create(
        Project(name=case.project_name, project_type=case.project_type)
    )

    imported_paths: list[Path] = list(case.source_files)
    if not imported_paths:
        imported_paths = materialize_inline_fallbacks(payload, scratch_dir)

    if not imported_paths:
        raise FileNotFoundError(
            f"Fixture {case.id} has no source files and no inline_docx/inline_xlsx fallback"
        )

    ingestion = IngestionService(session)
    for source_path in imported_paths:
        result = ingestion.import_file(project.id, source_path)
        if result.error:
            raise RuntimeError(f"Fixture import failed for {source_path.name}: {result.error}")
        if result.skipped:
            raise RuntimeError(f"Fixture import skipped for {source_path.name}")

    fact_repo = FactRepository(session)
    for fact in payload.get("facts", []):
        fact_repo.create(
            ProjectFact(
                project_id=project.id,
                key=str(fact["key"]),
                label=str(fact["label"]),
                value=str(fact["value"]),
                conflict_group=fact.get("conflict_group"),
            )
        )

    return case, project, imported_paths


def fixture_files_dir() -> Path:
    return _FILES_DIR
