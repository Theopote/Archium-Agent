"""Tests for slide recovery source parser."""

from __future__ import annotations

from pathlib import Path

import pytest

from archium.application.slide_recovery_source_parser import parse_source_page
from archium.domain.visual.render_scene import ImageNode, TextNode
from tests.fixtures.sample_files import create_sample_pdf


def test_parse_image_page(tmp_path: Path) -> None:
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow not installed")

    image_path = tmp_path / "page.png"
    Image.new("RGB", (1920, 1080), color=(255, 255, 255)).save(image_path)

    parsed = parse_source_page(image_path, workspace_dir=tmp_path / "ws")
    assert parsed.page_id == "page"
    assert parsed.source_kind == "image"
    assert parsed.preview_image_path is not None
    assert len(parsed.scene.nodes) == 1
    assert isinstance(parsed.scene.nodes[0], ImageNode)


def test_parse_pptx_slide(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100, 100, 400, 80).text = "测试标题"
    prs.save(str(pptx_path))

    parsed = parse_source_page(pptx_path, workspace_dir=tmp_path / "ws")
    assert parsed.page_id == "sample_slide_001"
    assert parsed.source_kind == "pptx"
    assert any(isinstance(node, TextNode) for node in parsed.scene.nodes)


def test_parse_pdf_page(tmp_path: Path) -> None:
    pytest.importorskip("fitz")
    pdf_path = tmp_path / "deck.pdf"
    create_sample_pdf(pdf_path, text="PDF 页面复活测试")

    parsed = parse_source_page(pdf_path, workspace_dir=tmp_path / "ws")
    assert parsed.source_kind == "pdf"
    assert parsed.page_id == "deck_page_001"
    assert parsed.preview_image_path is not None
    assert parsed.preview_image_path.is_file()
    assert len(parsed.scene.nodes) == 1
    assert isinstance(parsed.scene.nodes[0], ImageNode)
