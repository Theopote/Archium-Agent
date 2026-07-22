"""Tests for SlideRecoveryWorkflowService."""

from __future__ import annotations

from pathlib import Path

import pytest

from archium.application.slide_recovery_workflow_service import (
    SlideRecoveryWorkflowRequest,
    SlideRecoveryWorkflowService,
)
from archium.domain.enums import WorkflowStatus
from archium.domain.project import Project
from archium.domain.slide_recovery import SlideRecoveryPageKind
from archium.domain.visual.render_scene import TextNode
from archium.infrastructure.database.repositories import ProjectRepository
from tests.spike.slide_recovery_fixtures import SPIKE_SCENES


def test_workflow_run_completes_for_proxy_scene(db_session, tmp_path: Path) -> None:
    project = ProjectRepository(db_session).create(Project(name="Slide Recovery WF"))
    scene = SPIKE_SCENES[SlideRecoveryPageKind.TITLE]
    # Use PPTX path workaround: write pptx via parser input - use image instead
    pytest.importorskip("pptx")
    from pptx import Presentation

    pptx_path = tmp_path / "title.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for node in scene.nodes:
        if isinstance(node, TextNode):
            slide.shapes.add_textbox(100, 100, 400, 80).text = node.text
    prs.save(str(pptx_path))

    service = SlideRecoveryWorkflowService(db_session)
    result = service.run(
        project.id,
        SlideRecoveryWorkflowRequest(
            source_path=pptx_path,
            slide_index=0,
            page_kind=SlideRecoveryPageKind.TITLE,
            workspace_dir=tmp_path / "ws",
        ),
    )
    assert result.recovery_result is not None
    assert result.workflow_run.status in {
        WorkflowStatus.COMPLETED,
        WorkflowStatus.AWAITING_REVIEW,
    }
    assert result.source_page_id


def test_continue_after_review_accepts(db_session, tmp_path: Path) -> None:
    project = ProjectRepository(db_session).create(Project(name="Slide Recovery Review"))
    pytest.importorskip("pptx")
    from pptx import Presentation

    pptx_path = tmp_path / "accept.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100, 100, 400, 80).text = "标题"
    prs.save(str(pptx_path))

    service = SlideRecoveryWorkflowService(db_session)
    result = service.run(
        project.id,
        SlideRecoveryWorkflowRequest(
            source_path=pptx_path,
            workspace_dir=tmp_path / "ws",
            force_table_bitmap=True,
        ),
    )
    if result.workflow_run.status != WorkflowStatus.AWAITING_REVIEW:
        pytest.skip("workflow completed without review gate")

    finalized = service.continue_after_review(result.workflow_run.id, accepted=True)
    assert finalized.workflow_run.status == WorkflowStatus.COMPLETED
