"""Tests for ExportRoundTripService."""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock
from uuid import uuid4

import pytest

from archium.application.export_round_trip_service import (
    ExportRoundTripService,
    _derive_status,
    _drawing_integrity_checks,
    _text_recall,
)
from archium.application.visual.studio_scene_service import StudioSceneResult
from archium.config.settings import Settings
from archium.domain.export_round_trip import RoundTripStatus
from archium.domain.visual.render_scene import (
    BackgroundStyle,
    DrawingNode,
    RenderScene,
    TextNode,
)


def test_text_recall_partial_match() -> None:
    rate, missing = _text_recall(
        ("标题一", "正文段落"),
        ("标题一",),
    )
    assert rate == 0.5
    assert len(missing) == 1


def test_text_recall_full_match() -> None:
    rate, missing = _text_recall(
        ("标题一", "正文"),
        ("标题一", "正文"),
    )
    assert rate == 1.0
    assert not missing


def test_derive_status_blocked_on_blockers() -> None:
    status = _derive_status(
        text_match_rate=1.0,
        geometry_match_rate=1.0,
        similarity_score=0.95,
        drawing_issues=[],
        blockers=["页面数量不一致"],
        warnings=[],
    )
    assert status == RoundTripStatus.BLOCKED


def test_derive_status_needs_review_low_text_recall() -> None:
    status = _derive_status(
        text_match_rate=0.8,
        geometry_match_rate=1.0,
        similarity_score=0.95,
        drawing_issues=[],
        blockers=[],
        warnings=[],
    )
    assert status == RoundTripStatus.NEEDS_REVIEW


def test_drawing_cover_mode_is_integrity_issue() -> None:
    scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=1920,
        page_height=1080,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=[
            DrawingNode.model_construct(
                id="plan",
                node_type="drawing",
                x=100,
                y=100,
                width=800,
                height=600,
                z_index=1,
                storage_uri="asset://plan.png",
                fit_mode="cover",
            ),
        ],
    )
    issues = _drawing_integrity_checks(scene)
    assert any("fit_mode=cover" in issue for issue in issues)


def test_validate_pptx_ten_text_nodes_eight_found(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation as PptxPresentation

    presentation_id = uuid4()
    scene = RenderScene(
        slide_id=uuid4(),
        presentation_id=presentation_id,
        layout_plan_id=uuid4(),
        page_width=1920,
        page_height=1080,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=[
            TextNode(
                id=f"t{i}",
                x=100,
                y=100 + i * 50,
                width=400,
                height=40,
                z_index=i,
                text=f"节点文本{i}",
                font_family="Arial",
                font_size=18,
                color="#000",
                line_height=1.2,
            )
            for i in range(10)
        ],
    )

    pptx_path = tmp_path / "sparse.pptx"
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(8):
        slide.shapes.add_textbox(100, 100 + i * 50, 400, 40).text = f"节点文本{i}"
    prs.save(str(pptx_path))

    service = ExportRoundTripService.__new__(ExportRoundTripService)
    service._settings = Settings(_env_file=None)
    service._canvas = MagicMock()
    service._scene_service = MagicMock()
    service._scene_service.ensure_scenes_for_presentation.return_value = [
        StudioSceneResult(
            scene=scene,
            scene_hash="abc",
            preview_path=tmp_path / "preview.png",
            reused=False,
        )
    ]

    report = service.validate_pptx_export(
        presentation_id=presentation_id,
        pptx_path=pptx_path,
    )
    assert report.text_match_rate == pytest.approx(0.8, abs=0.05)
    assert report.status in {RoundTripStatus.NEEDS_REVIEW, RoundTripStatus.BLOCKED}


def test_validate_pptx_full_text_match(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation as PptxPresentation

    presentation_id = uuid4()
    scene = RenderScene(
        slide_id=uuid4(),
        presentation_id=presentation_id,
        layout_plan_id=uuid4(),
        page_width=1920,
        page_height=1080,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=[
            TextNode(
                id="t1",
                x=100,
                y=100,
                width=400,
                height=80,
                z_index=1,
                text="建筑汇报标题",
                font_family="Arial",
                font_size=24,
                color="#000",
                line_height=1.2,
            ),
        ],
    )

    pptx_path = tmp_path / "full.pptx"
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100, 100, 400, 80).text = "建筑汇报标题"
    prs.save(str(pptx_path))

    service = ExportRoundTripService.__new__(ExportRoundTripService)
    service._settings = Settings(_env_file=None)
    service._canvas = MagicMock()
    service._scene_service = MagicMock()
    service._scene_service.ensure_scenes_for_presentation.return_value = [
        StudioSceneResult(
            scene=scene,
            scene_hash="abc",
            preview_path=tmp_path / "preview.png",
            reused=False,
        )
    ]

    report = service.validate_pptx_export(
        presentation_id=presentation_id,
        pptx_path=pptx_path,
    )
    assert report.text_match_rate == 1.0
    assert report.status != RoundTripStatus.BLOCKED
