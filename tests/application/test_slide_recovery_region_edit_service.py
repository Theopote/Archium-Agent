"""Tests for slide recovery region editing."""

from __future__ import annotations

from pathlib import Path
from uuid import uuid4

import pytest
from archium.application.slide_recovery_region_edit_service import (
    SlideRecoveryRegionEditService,
    extract_regions,
    normalize_bbox,
    sanitize_region,
)
from archium.application.slide_recovery_service import SlideRecoveryRequest, SlideRecoveryService
from archium.application.slide_recovery_workflow_service import (
    SlideRecoveryWorkflowRequest,
    SlideRecoveryWorkflowService,
)
from archium.domain.enums import WorkflowStatus
from archium.domain.project import Project
from archium.domain.slide_recovery import (
    NormalizedBox,
    RecoveredPageRegion,
    SlideRecoveryPageKind,
)
from archium.domain.visual.render_scene import TextNode
from archium.infrastructure.database.repositories import ProjectRepository
from archium.infrastructure.slide_recovery.region_overlay_renderer import render_region_overlay
from tests.spike.slide_recovery_fixtures import SPIKE_SCENES


def test_normalize_bbox_clamps_to_page() -> None:
    box = normalize_bbox(x=0.95, y=0.9, width=0.2, height=0.2)
    assert box.x + box.width <= 1.0001
    assert box.y + box.height <= 1.0001


def test_render_region_overlay_writes_png(tmp_path: Path) -> None:
    pytest.importorskip("PIL")
    from PIL import Image

    source = tmp_path / "source.png"
    Image.new("RGB", (400, 225), color="#FFFFFF").save(source)
    regions = [
        RecoveredPageRegion(
            id=uuid4(),
            bbox=NormalizedBox(x=0.1, y=0.2, width=0.3, height=0.1),
            region_type="text",
            semantic_role="title",
            recovered_text="测试标题",
        )
    ]
    output = tmp_path / "overlay.png"
    render_region_overlay(source, regions, output)
    assert output.is_file()


def test_apply_region_edits_updates_workflow_state(db_session, tmp_path: Path) -> None:
    project = ProjectRepository(db_session).create(Project(name="Region Edit"))
    pytest.importorskip("pptx")
    from pptx import Presentation

    scene = SPIKE_SCENES[SlideRecoveryPageKind.TITLE]
    pptx_path = tmp_path / "title.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for node in scene.nodes:
        if isinstance(node, TextNode):
            slide.shapes.add_textbox(100, 100, 400, 80).text = node.text
    prs.save(str(pptx_path))

    workflow = SlideRecoveryWorkflowService(db_session)
    result = workflow.run(
        project.id,
        SlideRecoveryWorkflowRequest(
            source_path=pptx_path,
            slide_index=0,
            page_kind=SlideRecoveryPageKind.TITLE,
            workspace_dir=tmp_path / "ws",
        ),
    )
    regions = extract_regions(result)
    assert regions

    shifted = []
    for region in regions:
        bbox = region.bbox
        shifted.append(
            region.model_copy(
                update={
                    "bbox": bbox.model_copy(
                        update={"x": min(bbox.x + 0.02, 1.0 - bbox.width)}
                    )
                }
            )
        )

    editor = SlideRecoveryRegionEditService(db_session)
    updated = editor.apply_region_edits(result.workflow_run.id, shifted)

    assert updated.recovery_result is not None
    assert updated.workflow_run.status == WorkflowStatus.AWAITING_REVIEW
    assert updated.workflow_run.state.get("region_edit_applied") is True
    assert updated.hybrid_scene is not None
    assert len(updated.hybrid_scene.regions) == len(shifted)


def test_recover_page_with_precomputed_regions(db_session) -> None:
    scene = SPIKE_SCENES[SlideRecoveryPageKind.TITLE]
    regions = [
        RecoveredPageRegion(
            id=uuid4(),
            bbox=NormalizedBox(x=0.05, y=0.2, width=0.9, height=0.2),
            region_type="text",
            semantic_role="title",
            recovered_text="城市更新医疗综合体方案",
            source_node_id="title",
        )
    ]
    service = SlideRecoveryService(db_session)
    result = service.recover_page(
        SlideRecoveryRequest(
            source_page_id="manual",
            source_scene=scene,
            regions=regions,
            resolved_page_kind=SlideRecoveryPageKind.TITLE,
        )
    )
    assert result.hybrid_scene is not None
    title_nodes = [
        node
        for node in result.hybrid_scene.scene.nodes
        if getattr(node, "semantic_role", "") == "title"
    ]
    assert title_nodes
    assert sanitize_region(regions[0]).bbox.x == pytest.approx(0.05)
