"""Integration: publish template → compose slide → switch candidate."""

from __future__ import annotations

from pathlib import Path

from archium.application.visual.template_composition_service import TemplateCompositionService
from archium.application.visual.template_studio_service import TemplateStudioService
from archium.config.settings import Settings
from archium.domain.enums import ApprovalStatus, SlideType
from archium.domain.presentation import Presentation as DomainPresentation
from archium.domain.presentation import PresentationBrief, Storyline
from archium.domain.project import Project
from archium.domain.slide import SlideSpec
from archium.domain.visual.architectural_template import (
    TemplatePageType,
    TemplateSlot,
    TemplateSlotRole,
)
from archium.domain.visual.defaults import default_presentation_design_system
from archium.domain.visual.enums import DensityLevel, LayoutFamily, VisualContentType
from archium.domain.visual.visual_intent import VisualIntent
from archium.infrastructure.database.repositories import PresentationRepository, ProjectRepository
from archium.infrastructure.database.visual_repositories import (
    DesignSystemRepository,
    VisualIntentRepository,
)
from archium.ui.visual_service import apply_template_to_slide, select_layout_candidate
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session


def _write_pptx(path: Path) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.7))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "模板匹配封面"
    run.font.size = Pt(28)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.0), Inches(2.0))
    body.text_frame.text = "模板匹配正文区域"
    presentation.save(path)
    return path


def test_compose_slide_from_template_and_switch_candidate(
    db_session: Session,
    tmp_path: Path,
) -> None:
    settings = Settings(_env_file=None, output_path=tmp_path)
    project = ProjectRepository(db_session).create(Project(name="Template Compose"))
    presentations = PresentationRepository(db_session)
    presentation = presentations.create_presentation(
        DomainPresentation(project_id=project.id, title="Compose Deck")
    )
    brief = presentations.save_brief(
        PresentationBrief(
            project_id=project.id,
            presentation_id=presentation.id,
            title="Compose Deck",
            audience="甲方",
            purpose="模板匹配",
            core_message="模板可生成候选并切换。",
            approval_status=ApprovalStatus.APPROVED,
        )
    )
    storyline = presentations.save_storyline(
        Storyline(
            presentation_id=presentation.id,
            thesis="模板匹配验收。",
            approval_status=ApprovalStatus.APPROVED,
        )
    )
    presentation.current_brief_id = brief.id
    presentation.current_storyline_id = storyline.id
    presentations.update_presentation(presentation)

    DesignSystemRepository(db_session).save(default_presentation_design_system())
    slide = presentations.save_slide(
        SlideSpec(
            presentation_id=presentation.id,
            title="院区总平面",
            slide_type=SlideType.CONTENT,
            order=0,
            chapter_id="ch1",
            message="展示总平面与关键指标。",
            key_points=["北侧入口", "总建筑面积 10 万㎡"],
        )
    )
    intent = VisualIntentRepository(db_session).save(
        VisualIntent(
            slide_id=slide.id,
            presentation_id=presentation.id,
            communication_goal="突出场地关系",
            audience_takeaway="读懂总平面",
            visual_priority="drawing",
            dominant_content_type=VisualContentType.TEXT_ARGUMENT,
            preferred_layout_families=[LayoutFamily.TEXTUAL_ARGUMENT],
            density_level=DensityLevel.BALANCED,
        )
    )
    presentations.save_slide(slide.model_copy(update={"visual_intent_id": intent.id}))

    studio = TemplateStudioService(db_session, settings=settings)
    imported = studio.import_pptx(_write_pptx(tmp_path / "ref.pptx"), name="匹配模板")
    layout = imported.template.layouts[0]
    studio.update_page_type(imported.template.id, layout.id, TemplatePageType.TEXT_ARGUMENT)
    studio.upsert_slot(
        imported.template.id,
        layout.id,
        TemplateSlot(
            id="title",
            role=TemplateSlotRole.TITLE,
            x=0.8,
            y=0.4,
            width=8.0,
            height=0.7,
            auto_detected=False,
        ),
    )
    # Ensure a second layout candidate exists for switching.
    second = TemplateSlot(
        id="body",
        role=TemplateSlotRole.BODY,
        x=0.8,
        y=1.4,
        width=8.0,
        height=3.0,
        auto_detected=False,
    )
    template = studio.upsert_slot(imported.template.id, layout.id, second)
    # Duplicate page as second layout by updating slots on a cloned page via service internals:
    # Import already may have one page; create another layout by re-saving with two layouts.
    from archium.domain.visual.architectural_template import ArchitecturalTemplateLayout

    alt = ArchitecturalTemplateLayout(
        name="备选文字页",
        page_index=1,
        page_type=TemplatePageType.TEXT_ARGUMENT,
        slots=[
            TemplateSlot(
                id="alt_title",
                role=TemplateSlotRole.TITLE,
                x=1.0,
                y=0.5,
                width=7.5,
                height=0.7,
            ),
            TemplateSlot(
                id="alt_body",
                role=TemplateSlotRole.BODY,
                x=1.0,
                y=1.5,
                width=7.5,
                height=3.0,
            ),
        ],
        page_width=10,
        page_height=5.625,
        density_range=(0.2, 0.5),
    )
    template = studio._save(  # noqa: SLF001 — test seeds second layout
        template.model_copy(update={"layouts": [*template.layouts, alt]})
    )
    published = studio.publish(template.id)
    assert published.design_system_id is not None

    composed = apply_template_to_slide(
        db_session,
        slide_id=slide.id,
        template_id=published.id,
        candidate_count=3,
        settings=settings,
    )
    assert composed.layout_plan is not None
    assert composed.layout_plan.source_template_id == published.id
    assert len(composed.candidates) >= 1
    assert any(item.source_template_id == published.id for item in composed.candidates)

    if len(composed.candidates) >= 2:
        other = next(
            item for item in composed.candidates if item.id != composed.layout_plan.id
        )
        switched = select_layout_candidate(
            db_session,
            slide_id=slide.id,
            layout_plan_id=other.id,
        )
        assert switched.id == other.id
        assert switched.source_template_id == published.id

    # Composition service list should include the published template.
    listed = TemplateCompositionService(db_session, settings=settings).list_published_templates()
    assert any(item.id == published.id for item in listed)
