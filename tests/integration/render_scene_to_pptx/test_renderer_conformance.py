"""Integration tests for RenderScene → PPTX and renderer conformance."""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from archium.config.settings import Settings
from archium.infrastructure.renderers.pptx_renderer import PptxRenderer
from archium.infrastructure.renderers.pptxgen_cli import PptxGenCliRunner
from archium.infrastructure.renderers.renderer_conformance import assert_renderer_conformance
from pptx import Presentation

from tests.benchmark.architectural_slides.case_builders import build_benchmark_case
from tests.benchmark.architectural_slides.fixtures import ensure_case_assets
from tests.benchmark.architectural_slides.render_pipeline import (
    build_slide_content_bundle,
    compile_and_render_scene,
)


def _pptx_available() -> bool:
    runner = PptxGenCliRunner(Settings(_env_file=None))
    return (
        shutil.which("node") is not None
        and runner.is_available()
        and runner.layout_plan_script_path.exists()
    )


@pytest.mark.parametrize("case_id", ["case_001_site_plan", "case_006_project_hero"])
def test_renderer_conformance_scene_html_pptx(case_id: str, tmp_path: Path) -> None:
    if not _pptx_available():
        pytest.skip("Node/PptxGenJS unavailable")

    result = build_benchmark_case(case_id)
    case_dir = tmp_path / case_id
    case_dir.mkdir()
    ensure_case_assets(case_id, case_dir / "assets")
    scene_render = compile_and_render_scene(result, case_dir)

    pptx_path = case_dir / "output.pptx"
    PptxRenderer(Settings(_env_file=None)).export_pptx(
        scene_render.scene,
        pptx_path,
        title=result.slide.title,
    )
    assert pptx_path.is_file()
    assert pptx_path.stat().st_size > 500

    report = assert_renderer_conformance(scene_render.scene, pptx_path=pptx_path)
    assert report.passed, report.issues

    presentation = Presentation(pptx_path)
    assert len(presentation.slides) == 1
    slide_text = " ".join(
        shape.text.strip()
        for shape in presentation.slides[0].shapes
        if getattr(shape, "text", "") and shape.text.strip()
    )
    assert result.slide.title in slide_text


def test_pptx_contains_editable_text_not_flattened(tmp_path: Path) -> None:
    if not _pptx_available():
        pytest.skip("Node/PptxGenJS unavailable")

    result = build_benchmark_case("case_001_site_plan")
    case_dir = tmp_path / "case"
    case_dir.mkdir()
    ensure_case_assets("case_001_site_plan", case_dir / "assets")
    build = build_slide_content_bundle(result.plan, case_dir / "assets", result.slide)
    from archium.application.visual.render_scene_compiler import RenderSceneCompiler

    scene = RenderSceneCompiler().compile(
        slide=result.slide,
        layout_plan=result.plan,
        design_system=result.design_system,
        content_bundle=build.bundle,
        visual_intent=result.intent,
    )
    pptx_path = case_dir / "output.pptx"
    PptxRenderer(Settings(_env_file=None)).export_pptx(
        scene,
        pptx_path,
        title=result.slide.title,
    )

    presentation = Presentation(pptx_path)
    slide = presentation.slides[0]
    text_shapes = [
        shape for shape in slide.shapes if getattr(shape, "text", "") and shape.text.strip()
    ]
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
    assert len(text_shapes) >= 2
    assert len(picture_shapes) >= 1
    assert len(slide.shapes) > 2
