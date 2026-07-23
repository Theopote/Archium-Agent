"""Smoke: STRUCTURED PPTX export emits real masters/layouts and passes OOXML checks."""

from __future__ import annotations

import shutil
from pathlib import Path
from uuid import uuid4

import pytest
from archium.config.settings import Settings
from archium.domain.visual.pptx_structure import PlaceholderKind, PptxStructureMode
from archium.infrastructure.renderers.pptx_structure_catalog import p0_structured_spike_spec
from archium.domain.visual.render_scene import BackgroundStyle, RenderScene, TextNode
from archium.infrastructure.renderers.pptx_ooxml_structure import (
    inspect_pptx_ooxml_structure,
    require_structured_ooxml,
)
from archium.infrastructure.renderers.pptx_renderer import PptxRenderer
from archium.infrastructure.renderers.pptxgen_cli import PptxGenCliRunner
from pptx import Presentation
from tests.smoke.artifact_publish import publish_smoke_artifact

pytestmark = pytest.mark.smoke


def _runner_available() -> bool:
    runner = PptxGenCliRunner(Settings(_env_file=None))
    return runner.is_available() and runner.layout_plan_script_path.exists()


@pytest.mark.skipif(shutil.which("node") is None, reason="Node.js not installed")
def test_structured_pptx_emits_masters_layouts_and_placeholders(tmp_path: Path) -> None:
    if not _runner_available():
        pytest.skip(
            "PptxGenJS runtime unavailable — run npm install in "
            "archium/infrastructure/renderers/pptxgen"
        )

    title_scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=10.0,
        page_height=5.625,
        background=BackgroundStyle(color="#FFFFFF"),
        source_layout_family="hero",
        nodes=[
            TextNode(
                id="title",
                x=0.5,
                y=1.8,
                width=9.0,
                height=1.0,
                text="Structured Title",
                font_family="Arial",
                font_size=32,
                color="#111111",
                line_height=1.2,
                semantic_role="title",
            ),
            TextNode(
                id="subtitle",
                x=0.5,
                y=3.0,
                width=9.0,
                height=0.6,
                text="Native master/layout export",
                font_family="Arial",
                font_size=18,
                color="#444444",
                line_height=1.2,
                semantic_role="subtitle",
            ),
        ],
    )
    content_scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=10.0,
        page_height=5.625,
        background=BackgroundStyle(color="#FFFFFF"),
        source_layout_family="textual_argument",
        nodes=[
            TextNode(
                id="title",
                x=0.5,
                y=0.4,
                width=9.0,
                height=0.7,
                text="Content Layout",
                font_family="Arial",
                font_size=24,
                color="#111111",
                line_height=1.2,
                semantic_role="title",
            ),
            TextNode(
                id="body",
                x=0.5,
                y=1.3,
                width=9.0,
                height=3.5,
                text="Body text bound through a placeholder.",
                font_family="Arial",
                font_size=16,
                color="#222222",
                line_height=1.3,
                semantic_role="body",
            ),
        ],
    )
    picture_scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=10.0,
        page_height=5.625,
        background=BackgroundStyle(color="#FFFFFF"),
        source_layout_family="drawing_focus",
        nodes=[
            TextNode(
                id="title",
                x=0.5,
                y=0.4,
                width=9.0,
                height=0.7,
                text="Picture Layout",
                font_family="Arial",
                font_size=24,
                color="#111111",
                line_height=1.2,
                semantic_role="title",
            ),
            TextNode(
                id="caption",
                x=7.0,
                y=1.3,
                width=2.5,
                height=3.5,
                text="Caption beside picture placeholder.",
                font_family="Arial",
                font_size=14,
                color="#333333",
                line_height=1.2,
                semantic_role="caption",
            ),
        ],
    )

    spike = p0_structured_spike_spec()
    assert len(spike.masters) == 1
    assert len(spike.layouts) == 3
    kinds = {ph.placeholder_type for layout in spike.layouts for ph in layout.placeholder_specs}
    assert PlaceholderKind.TITLE in kinds
    assert PlaceholderKind.BODY in kinds
    assert PlaceholderKind.IMAGE in kinds
    assert PlaceholderKind.SLIDE_NUMBER in kinds

    output_path = (tmp_path / "structured.editable.pptx").resolve()
    renderer = PptxRenderer(Settings(_env_file=None))
    rendered = renderer.export_presentation(
        title="Archium Structured Structure Smoke",
        scenes=[
            (title_scene, "title notes"),
            (content_scene, "body notes"),
            (picture_scene, "picture notes"),
        ],
        output_path=output_path,
        structure_mode=PptxStructureMode.STRUCTURED,
        structure=spike,
        validate_ooxml=True,
    )

    assert rendered.exists()
    report = inspect_pptx_ooxml_structure(rendered)
    assert report.valid
    assert report.structure_mode == PptxStructureMode.STRUCTURED
    # P0-5: one Master, three Layouts after expander consolidation.
    assert len(report.master_parts) == 1
    assert len(report.layout_parts) == 3
    assert report.slide_to_layout
    assert report.layout_to_master
    assert report.placeholder_count >= 1
    require_structured_ooxml(rendered)

    presentation = Presentation(rendered)
    assert len(presentation.slides) == 3
    assert presentation.slides[0].slide_layout is not None
    assert presentation.slides[0].slide_layout.slide_master is not None

    publish_smoke_artifact(rendered, "pptxgen_structured_structure.editable.pptx")
