"""Windows real-machine smoke — PptxGen, checkpoint, Unicode paths, basic export.

Run on ``windows-latest`` via ``.github/workflows/windows-smoke.yml`` (nightly + manual RC).
Skipped on non-Windows hosts so default Linux CI stays unchanged.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from archium.application.ingestion_service import IngestionService
from archium.application.presentation_models import PresentationRequest
from archium.application.presentation_workflow_service import PresentationWorkflowService
from archium.config.settings import Settings
from archium.domain.enums import ProjectType, WorkflowStatus
from archium.domain.project import Project
from archium.infrastructure.database.repositories import DocumentRepository, ProjectRepository
from archium.infrastructure.llm import MockLLMProvider
from archium.infrastructure.renderers.pptxgen_cli import PptxGenCliRunner
from archium.workflow.checkpointer import WorkflowCheckpointerManager
from pptx import Presentation
from sqlalchemy.orm import Session
from tests.fixtures.mock_llm import pipeline_mock_selector
from tests.golden.fixtures.loader import seed_fixture_case
from tests.golden.regression.loader import load_regression_case, seed_regression_case

pytestmark = [
    pytest.mark.smoke,
    pytest.mark.windows_smoke,
    pytest.mark.skipif(sys.platform != "win32", reason="Windows platform smoke (see windows-smoke.yml)"),
]

_SPEC_PATH = Path(__file__).resolve().parents[1] / "fixtures" / "pptxgen" / "smoke.spec.json"
_CASE_A = Path(__file__).resolve().parents[1] / "golden" / "regression" / "cases" / "case_a_hospital.json"
_CASE_E_MANIFEST = (
    Path(__file__).resolve().parents[1] / "golden" / "fixtures" / "manifests" / "case_e_real_paths.fixture.json"
)


def _pptxgen_runner() -> PptxGenCliRunner:
    runner = PptxGenCliRunner(Settings(_env_file=None))
    if not runner.is_available():
        pytest.skip("PptxGenJS runtime unavailable — run npm ci in archium/infrastructure/renderers/pptxgen")
    return runner


def test_pptxgen_render_and_reopen_on_windows(tmp_path: Path) -> None:
    runner = _pptxgen_runner()
    output_path = (tmp_path / "smoke.editable.pptx").resolve()
    rendered = runner.render(_SPEC_PATH.resolve(), output_path)

    assert rendered.exists()
    assert rendered.stat().st_size > 500
    presentation = Presentation(rendered)
    assert len(presentation.slides) == 2
    first_notes = presentation.slides[0].notes_slide.notes_text_frame.text
    assert "Smoke test speaker note" in first_notes


def test_pptxgen_export_to_unicode_output_path(tmp_path: Path) -> None:
    """Export PPTX to a Chinese + spaced directory — primary Windows delivery path."""
    runner = _pptxgen_runner()
    output_dir = tmp_path / "中文输出目录" / "项目 资料"
    output_dir.mkdir(parents=True)
    output_path = (output_dir / "汇报 文件.pptx").resolve()

    rendered = runner.render(_SPEC_PATH.resolve(), output_path)
    assert rendered.exists()
    assert rendered.stat().st_size > 500

    presentation = Presentation(rendered)
    assert len(presentation.slides) == 2


def test_checkpoint_db_lifecycle_on_windows(tmp_path: Path) -> None:
    db_path = tmp_path / "checkpoints" / "workflow checkpoints.db"
    with WorkflowCheckpointerManager(db_path) as manager:
        _ = manager.saver
        assert db_path.exists()

    db_path.unlink()
    assert not db_path.exists()


def test_workflow_checkpoint_cleanup_after_run(
    db_session: Session,
    test_settings: Settings,
) -> None:
    checkpoint_path = test_settings.workflow_checkpoint_path
    manager = WorkflowCheckpointerManager(checkpoint_path)
    project = ProjectRepository(db_session).create(
        Project(name="Windows Checkpoint", project_type=ProjectType.HEALTHCARE)
    )
    service = PresentationWorkflowService(
        db_session,
        MockLLMProvider(selector=pipeline_mock_selector),
        settings=test_settings,
        checkpointer_manager=manager,
    )
    payload = PresentationRequest(
        title="Windows Checkpoint",
        audience="QA",
        purpose="Verify SQLite checkpoint cleanup on Windows",
        duration_minutes=10,
        target_slide_count=2,
        core_message="Close checkpoint DB after workflow",
    )
    try:
        result = service.run(project.id, payload, export_marp=False)
        assert result.workflow_run.status == WorkflowStatus.COMPLETED
        assert checkpoint_path.exists()
    finally:
        service.close()

    checkpoint_path.unlink()
    assert not checkpoint_path.exists()


def test_unicode_and_spaced_path_ingestion(tmp_path: Path, db_session: Session) -> None:
    case, project, imported_paths = seed_fixture_case(db_session, _CASE_E_MANIFEST, scratch_dir=tmp_path)

    assert case.id == "case_e_real_paths"
    assert any(ord(char) > 127 for char in str(imported_paths[0]))
    assert any(" " in str(path) for path in imported_paths)
    assert len(imported_paths) >= 4

    chunks = DocumentRepository(db_session).list_chunks_by_project(project.id)
    assert len(chunks) >= 3

    extensions = {path.suffix.lower() for path in imported_paths}
    assert extensions >= {".docx", ".pdf", ".pptx", ".jpg"}


def test_low_resolution_image_on_unicode_path(db_session: Session, tmp_path: Path) -> None:
    project = ProjectRepository(db_session).create(
        Project(name="低清图片测试", project_type=ProjectType.HEALTHCARE)
    )
    from tests.golden.fixtures.loader import materialize_inline_image

    image_path = materialize_inline_image(
        tmp_path / "中文路径" / "现场 低清.jpg",
        {"width": 32, "height": 24},
    )
    result = IngestionService(db_session).import_file(project.id, image_path)
    assert not result.error
    assert not result.skipped

    chunks = DocumentRepository(db_session).list_chunks_by_project(project.id)
    assert len(chunks) >= 1


def test_golden_case_a_json_export_on_windows(
    db_session: Session,
    test_settings: Settings,
) -> None:
    case = load_regression_case(_CASE_A)
    _, project = seed_regression_case(db_session, _CASE_A)
    service = PresentationWorkflowService(
        db_session,
        MockLLMProvider(selector=pipeline_mock_selector),
        settings=test_settings,
    )
    try:
        result = service.run(
            project.id,
            case.request,
            export_json=True,
            export_marp=False,
            require_brief_review=False,
            require_storyline_review=False,
            require_slides_review=False,
        )
    finally:
        service.close()

    assert result.workflow_run.status == WorkflowStatus.COMPLETED
    assert result.render.json_path is not None
    payload = json.loads(result.render.json_path.read_text(encoding="utf-8"))
    assert len(payload.get("slides", [])) >= case.expectations.get("min_slides", 1)
