"""Smoke: NATIVE_DATA_BACKED chart/table emit real OOXML chart/table parts."""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from uuid import uuid4

import pytest
from archium.config.settings import Settings
from archium.domain.export_fidelity import ChartExportMode
from archium.domain.visual.render_scene import (
    BackgroundStyle,
    ChartNode,
    ChartSeriesData,
    RenderScene,
    TableNode,
)
from archium.infrastructure.renderers.pptx_renderer import PptxRenderer
from archium.infrastructure.renderers.pptxgen_cli import PptxGenCliRunner
from pptx import Presentation
from tests.smoke.artifact_publish import publish_smoke_artifact

pytestmark = pytest.mark.smoke


def _runner_available() -> bool:
    runner = PptxGenCliRunner(Settings(_env_file=None))
    return runner.is_available() and runner.layout_plan_script_path.exists()


@pytest.mark.skipif(shutil.which("node") is None, reason="Node.js not installed")
def test_native_chart_and_table_export_contains_ooxml_parts(tmp_path: Path) -> None:
    if not _runner_available():
        pytest.skip(
            "PptxGenJS runtime unavailable — run npm install in "
            "archium/infrastructure/renderers/pptxgen"
        )

    scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=10.0,
        page_height=5.625,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=[
            ChartNode(
                id="chart1",
                x=0.4,
                y=0.8,
                width=4.6,
                height=3.8,
                chart_type="bar",
                title="指标",
                series=[
                    ChartSeriesData(
                        name="用量",
                        labels=["一期", "二期", "三期"],
                        values=[12.0, 18.5, 9.0],
                    )
                ],
            ),
            TableNode(
                id="table1",
                x=5.2,
                y=0.8,
                width=4.4,
                height=3.8,
                headers=["分期", "面积"],
                rows=[["一期", "12000"], ["二期", "9800"], ["三期", "7600"]],
            ),
        ],
    )

    output = (tmp_path / "native_chart_table.pptx").resolve()
    rendered = PptxRenderer(Settings(_env_file=None)).export_presentation(
        title="Native Chart Table Smoke",
        scenes=[(scene, "notes")],
        output_path=output,
        chart_export_mode=ChartExportMode.NATIVE_DATA_BACKED,
        validate_ooxml=False,
    )
    assert rendered.exists()

    with zipfile.ZipFile(rendered) as archive:
        names = set(archive.namelist())
    assert any(name.startswith("ppt/charts/chart") for name in names)
    # Embedded workbook backing the chart data.
    assert any("embeddings" in name.lower() or name.endswith(".xlsx") for name in names)
    # Native table lives on the slide as a graphic frame; python-pptx can see it.
    presentation = Presentation(rendered)
    assert len(presentation.slides) == 1
    has_table = any(shape.has_table for shape in presentation.slides[0].shapes)
    assert has_table

    publish_smoke_artifact(rendered, "pptxgen_native_chart_table.editable.pptx")
