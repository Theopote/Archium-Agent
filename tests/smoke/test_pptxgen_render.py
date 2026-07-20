"""PptxGenJS end-to-end smoke test — real Node render + python-pptx verification."""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from archium.config.settings import Settings
from archium.infrastructure.renderers.pptxgen_cli import PptxGenCliRunner
from pptx import Presentation
from tests.smoke.artifact_publish import publish_smoke_artifact

pytestmark = pytest.mark.smoke

_SPEC_PATH = Path(__file__).resolve().parents[1] / "fixtures" / "pptxgen" / "smoke.spec.json"
_PPTXGEN_DIR = (
    Path(__file__).resolve().parents[2]
    / "archium"
    / "infrastructure"
    / "renderers"
    / "pptxgen"
)


def _runner_available() -> bool:
    runner = PptxGenCliRunner(Settings(_env_file=None))
    return runner.is_available()


def _slide_texts(slide: object) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        text = getattr(shape, "text", "")
        if text and str(text).strip():
            texts.append(str(text).strip())
    return texts


@pytest.mark.skipif(shutil.which("node") is None, reason="Node.js not installed")
def test_pptxgen_smoke_render_and_reopen(tmp_path: Path) -> None:
    if not _runner_available():
        pytest.skip("PptxGenJS runtime unavailable — run npm install in archium/infrastructure/renderers/pptxgen")

    output_path = (tmp_path / "smoke.editable.pptx").resolve()
    runner = PptxGenCliRunner(Settings(_env_file=None))
    rendered = runner.render(_SPEC_PATH.resolve(), output_path)

    assert rendered.exists()
    assert rendered.stat().st_size > 500

    presentation = Presentation(rendered)
    assert len(presentation.slides) == 2

    first_text = " ".join(_slide_texts(presentation.slides[0]))
    assert "Archium PptxGen Smoke" in first_text

    first_notes = presentation.slides[0].notes_slide.notes_text_frame.text
    assert "Smoke test speaker note" in first_notes

    second_text = " ".join(_slide_texts(presentation.slides[1]))
    assert "要点页" in second_text
    second_notes = presentation.slides[1].notes_slide.notes_text_frame.text
    assert "中文路径" in second_notes

    publish_smoke_artifact(rendered, "pptxgen_smoke.editable.pptx")


@pytest.mark.skipif(shutil.which("node") is None, reason="Node.js not installed")
def test_pptxgen_node_modules_present() -> None:
    node_modules = _PPTXGEN_DIR / "node_modules" / "pptxgenjs"
    if not node_modules.exists():
        pytest.skip("pptxgenjs not installed — run npm install in pptxgen directory")
    assert (_PPTXGEN_DIR / "render.mjs").exists()
