"""Smoke test: LayoutPlan instruction deck → real PptxGenJS PPTX."""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from uuid import uuid4

import pytest
from archium.config.settings import Settings
from archium.domain.visual import (
    LayoutElement,
    LayoutElementRole,
    LayoutFamily,
    LayoutPlan,
    default_presentation_design_system,
)
from archium.domain.visual.enums import LayoutContentType
from archium.infrastructure.renderers.pptxgen.layout_plan_adapter import (
    PptxLayoutPlanAdapter,
    SlideContentBundle,
)
from archium.infrastructure.renderers.pptxgen_cli import PptxGenCliRunner
from pptx import Presentation

pytestmark = pytest.mark.smoke

_ARTIFACT_DIR = Path(__file__).resolve().parent / "artifacts"
_PPTXGEN_DIR = (
    Path(__file__).resolve().parents[2]
    / "archium"
    / "infrastructure"
    / "renderers"
    / "pptxgen"
)


def _runner_available() -> bool:
    runner = PptxGenCliRunner(Settings(_env_file=None))
    return runner.is_available() and runner.layout_plan_script_path.exists()


def _slide_texts(slide: object) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        text = getattr(shape, "text", "")
        if text and str(text).strip():
            texts.append(str(text).strip())
    return texts


@pytest.mark.skipif(shutil.which("node") is None, reason="Node.js not installed")
def test_layout_plan_smoke_render_preserves_text(tmp_path: Path) -> None:
    if not _runner_available():
        pytest.skip(
            "PptxGenJS runtime unavailable — run npm install in "
            "archium/infrastructure/renderers/pptxgen"
        )

    design = default_presentation_design_system()
    plan = LayoutPlan(
        slide_id=uuid4(),
        layout_family=LayoutFamily.DRAWING_FOCUS,
        layout_variant="full_canvas",
        page_width=10.0,
        page_height=5.625,
        design_system_id=design.id,
        visual_intent_id=uuid4(),
        reading_order=["title", "body"],
        elements=[
            LayoutElement(
                id="title",
                role=LayoutElementRole.TITLE,
                content_type=LayoutContentType.TEXT,
                text_content="LayoutPlan 执行渲染",
                x=0.7,
                y=0.45,
                width=8.6,
                height=0.5,
                style_token="title",
            ),
            LayoutElement(
                id="body",
                role=LayoutElementRole.BODY_TEXT,
                content_type=LayoutContentType.TEXT,
                text_content="坐标由 LayoutPlan 决定，renderer 不得重排。",
                x=0.7,
                y=1.2,
                width=8.6,
                height=1.5,
                style_token="body",
            ),
        ],
    )
    deck = PptxLayoutPlanAdapter().render_deck(
        title="LayoutPlan Smoke",
        slides=[(plan, design, SlideContentBundle(page_number=1))],
    )
    deck_path = tmp_path / "presentation.layout_instructions.json"
    deck_path.write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
    output_path = (tmp_path / "layout_plan.editable.pptx").resolve()

    runner = PptxGenCliRunner(Settings(_env_file=None))
    rendered = runner.render_layout_instructions(deck_path.resolve(), output_path)

    assert rendered.exists()
    assert rendered.stat().st_size > 500

    presentation = Presentation(rendered)
    assert len(presentation.slides) == 1
    text = " ".join(_slide_texts(presentation.slides[0]))
    assert "LayoutPlan 执行渲染" in text
    assert "坐标由 LayoutPlan 决定" in text

    # Spot-check that title shape left/top roughly match instruction inches.
    # python-pptx uses EMUs (914400 per inch).
    title_shape = next(
        shape
        for shape in presentation.slides[0].shapes
        if getattr(shape, "text", "") and "LayoutPlan 执行渲染" in shape.text
    )
    assert abs(title_shape.left / 914400 - 0.7) < 0.05
    assert abs(title_shape.top / 914400 - 0.45) < 0.05

    _ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    artifact = _ARTIFACT_DIR / "layout_plan_smoke.editable.pptx"
    artifact.write_bytes(rendered.read_bytes())
    assert (_PPTXGEN_DIR / "render-plan.mjs").exists()
