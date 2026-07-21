"""Phase 2 parser fixes: assets, groups, placeholders, page area, repeat signatures."""

from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from archium.application.visual.asset_path_resolver import is_machine_absolute_path
from archium.domain.visual.reference_slide import (
    REFERENCE_TEMPLATE_ASSET_ORIGIN,
    ReferenceElement,
    ReferenceElementType,
)
from archium.infrastructure.template.reference_pptx_parser import (
    ReferencePptxParser,
    _infer_element_type,
    _mark_repeated_elements,
    _structural_signature,
)


def _png(path: Path, color: tuple[int, int, int] = (200, 200, 200)) -> Path:
    Image.new("RGB", (320, 240), color).save(path)
    return path


def test_picture_blob_extracted_to_workspace_assets(tmp_path: Path) -> None:
    png = _png(tmp_path / "plan.png")
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    pic = slide.shapes.add_picture(str(png), Inches(0.5), Inches(1.0), width=Inches(4.0))
    pic.name = "总平面图"
    pptx = tmp_path / "pic.pptx"
    presentation.save(pptx)

    workspace = tmp_path / "ws"
    parsed = ReferencePptxParser().parse(pptx, workspace_dir=workspace, capture_screenshots=False)
    assets = parsed.slides[0].image_assets
    assert assets
    asset = assets[0]
    assert asset.asset_origin == REFERENCE_TEMPLATE_ASSET_ORIGIN
    assert asset.relative_path.startswith("assets/slide_001/image_")
    assert asset.content_hash
    assert not is_machine_absolute_path(asset.relative_path)
    disk = workspace / asset.relative_path
    assert disk.is_file()
    assert disk.stat().st_size > 50
    element = next(
        e
        for e in parsed.slides[0].iter_elements()
        if e.element_type in {ReferenceElementType.IMAGE, ReferenceElementType.DRAWING}
    )
    assert element.asset_id == asset.id


def test_group_children_are_recursively_parsed_with_workspace(tmp_path: Path) -> None:
    parser = ReferencePptxParser()
    child_picture = SimpleNamespace(
        left=Inches(1).emu,
        top=Inches(1).emu,
        width=Inches(2).emu,
        height=Inches(2).emu,
        shape_type=MSO_SHAPE_TYPE.PICTURE,
        has_text_frame=False,
        name="ChildPicture",
        is_placeholder=False,
        image=SimpleNamespace(blob=b"\x89PNG\r\n\x1a\n" + b"\x00" * 40, content_type="image/png"),
        _element=SimpleNamespace(nvPicPr=SimpleNamespace(cNvPr={"descr": "", "title": ""})),
        fill=None,
    )
    child_text = SimpleNamespace(
        left=Inches(1).emu,
        top=Inches(3.2).emu,
        width=Inches(2).emu,
        height=Inches(0.4).emu,
        shape_type=MSO_SHAPE_TYPE.TEXT_BOX,
        has_text_frame=True,
        name="ChildCaption",
        is_placeholder=False,
        text_frame=SimpleNamespace(
            paragraphs=[
                SimpleNamespace(
                    text="图注：总平面",
                    runs=[
                        SimpleNamespace(
                            font=SimpleNamespace(name="Arial", size=None, color=None)
                        )
                    ],
                )
            ]
        ),
        fill=None,
    )
    group = SimpleNamespace(
        left=Inches(0.5).emu,
        top=Inches(0.5).emu,
        width=Inches(4).emu,
        height=Inches(4).emu,
        shape_type=MSO_SHAPE_TYPE.GROUP,
        has_text_frame=False,
        name="CardGroup",
        is_placeholder=False,
        shapes=[child_picture, child_text],
        fill=None,
    )
    element, assets, texts, _, _ = parser._parse_shape(
        group,
        slide_index=0,
        z_index=0,
        id_prefix="s001_e001",
        page_width=10.0,
        page_height=5.625,
        workspace=tmp_path,
        layout_name="Blank",
        master_name="Master",
        image_seq={"n": 0},
    )
    assert element is not None
    assert element.element_type == ReferenceElementType.GROUP
    assert len(element.children) == 2
    assert any(c.element_type == ReferenceElementType.IMAGE for c in element.children)
    assert any(c.element_type == ReferenceElementType.TEXT for c in element.children)
    assert assets
    assert assets[0].relative_path.startswith("assets/slide_001/")
    assert "总平面" in "".join(texts)


def test_placeholder_shapes_are_classified(tmp_path: Path) -> None:
    presentation = Presentation()
    # Title layout includes native placeholders.
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    pptx = tmp_path / "ph.pptx"
    presentation.save(pptx)
    parsed = ReferencePptxParser().parse(
        pptx, workspace_dir=tmp_path / "ws", capture_screenshots=False
    )
    types = {e.element_type for e in parsed.slides[0].iter_elements()}
    assert ReferenceElementType.PLACEHOLDER in types


def test_decoration_uses_actual_page_width_not_hardcoded_10() -> None:
    """Large-area decoration threshold must scale with real page_width * page_height."""
    shape = SimpleNamespace(is_placeholder=False)
    # has_text=True with empty text avoids the blanket "not has_text → decoration" branch.
    on_narrow_page = _infer_element_type(
        shape=shape,
        shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE,
        has_text=True,
        text="",
        is_picture=False,
        width=4.0,
        height=3.0,  # area 12; page 5x5=25 → ratio 0.48 ≥ 0.175
        page_width=5.0,
        page_height=5.0,
    )
    on_wide_page = _infer_element_type(
        shape=shape,
        shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE,
        has_text=True,
        text="",
        is_picture=False,
        width=4.0,
        height=3.0,  # area 12; page 20x10=200 → ratio 0.06 < 0.175
        page_width=20.0,
        page_height=10.0,
    )
    assert on_narrow_page == ReferenceElementType.DECORATION
    assert on_wide_page == ReferenceElementType.SHAPE


def test_repeat_signature_ignores_generic_textbox_names() -> None:
    def make(text: str, name: str, slide_index: int) -> ReferenceElement:
        sig = _structural_signature(
            element_type=ReferenceElementType.TEXT,
            x=0.7,
            y=1.5,
            width=4.0,
            height=1.0,
            font_name="Arial",
            fill_color=None,
            text=text,
            asset_hash="",
            layout_name="Blank",
            master_name="Master",
            placeholder=False,
        )
        return ReferenceElement(
            id=f"s{slide_index}_e1",
            element_type=ReferenceElementType.TEXT,
            x=0.7,
            y=1.5,
            width=4.0,
            height=1.0,
            text=text,
            source_shape_name=name,
            structural_signature=sig,
            semantic_role="body",
        )

    from archium.domain.visual.reference_slide import ReferenceSlideSnapshot

    slides = []
    for i in range(3):
        # Same auto-name "TextBox 1" but unique body text → different signatures.
        slides.append(
            ReferenceSlideSnapshot(
                slide_index=i,
                slide_id=f"slide_{i + 1:03d}",
                elements=[make(f"unique body {i}", "TextBox 1", i)],
            )
        )
    _mark_repeated_elements(slides)
    assert all(not e.repeats_across_pages for s in slides for e in s.elements)
    assert all(not e.likely_background_or_decoration for s in slides for e in s.elements)

    # Identical chrome footer on every page should repeat.
    chrome_slides = []
    for i in range(3):
        chrome_slides.append(
            ReferenceSlideSnapshot(
                slide_index=i,
                slide_id=f"slide_{i + 1:03d}",
                elements=[make("页脚 · 机密", "TextBox 1", i)],
            )
        )
    _mark_repeated_elements(chrome_slides)
    assert all(e.repeats_across_pages for s in chrome_slides for e in s.elements)
