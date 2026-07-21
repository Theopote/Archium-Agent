"""Parser tests: real Picture shapes (PNG) must be inferable as drawings."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from archium.domain.visual.reference_slide import ReferenceElementType
from archium.infrastructure.template.reference_pptx_parser import ReferencePptxParser


def _make_png(path: Path, color: tuple[int, int, int] = (220, 220, 220)) -> Path:
    Image.new("RGB", (640, 480), color).save(path)
    return path


def _write_drawing_picture_pptx(path: Path, png: Path) -> Path:
    """Title + real Picture (no text frame) + caption — classic site-plan page."""
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)

    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.6))
    title.text_frame.text = "总平面结构与关键指标"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    picture = slide.shapes.add_picture(str(png), Inches(0.6), Inches(1.1), width=Inches(5.5))
    picture.name = "总平面图"
    picture._element.nvPicPr.cNvPr.set("descr", "院区总平面图纸")

    caption = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(5.5), Inches(0.4))
    caption.text_frame.text = "图：总平面图"
    caption.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # Control: a photo picture with photo cues should stay IMAGE.
    photo_png = png.parent / "photo.png"
    _make_png(photo_png, (180, 160, 140))
    slide2 = presentation.slides.add_slide(blank)
    t2 = slide2.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.6))
    t2.text_frame.text = "入口现场照片"
    t2.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    pic2 = slide2.shapes.add_picture(str(photo_png), Inches(0.6), Inches(1.1), width=Inches(4.5))
    pic2.name = "现场照片_01"
    pic2._element.nvPicPr.cNvPr.set("descr", "现场照片")
    c2 = slide2.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(4.5), Inches(0.4))
    c2.text_frame.text = "现场照片：车行冲突"

    presentation.save(path)
    return path


def test_real_picture_without_text_frame_inferred_as_drawing(tmp_path: Path) -> None:
    png = _make_png(tmp_path / "plan.png")
    pptx = _write_drawing_picture_pptx(tmp_path / "drawing.pptx", png)
    presentation = ReferencePptxParser().parse(
        pptx, workspace_dir=tmp_path / "ws", capture_screenshots=False
    )
    assert len(presentation.slides) == 2

    plan_slide = presentation.slides[0]
    pictures = [
        e
        for e in plan_slide.elements
        if e.element_type in {ReferenceElementType.IMAGE, ReferenceElementType.DRAWING}
    ]
    assert pictures, "expected a picture element"
    assert pictures[0].text == ""  # real pictures have empty text frames
    assert pictures[0].element_type == ReferenceElementType.DRAWING
    assert pictures[0].semantic_role == "drawing"
    assert pictures[0].alt_text
    assert any("drawing_inference" in n for n in pictures[0].style_notes)

    photo_slide = presentation.slides[1]
    photos = [
        e
        for e in photo_slide.elements
        if e.element_type in {ReferenceElementType.IMAGE, ReferenceElementType.DRAWING}
    ]
    assert photos
    assert photos[0].element_type == ReferenceElementType.IMAGE


def test_picture_named_floor_plan_without_neighbors(tmp_path: Path) -> None:
    png = _make_png(tmp_path / "floor.png")
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    pic = slide.shapes.add_picture(str(png), Inches(1.0), Inches(1.0), width=Inches(6.0))
    pic.name = "FloorPlan_L1"
    pic._element.nvPicPr.cNvPr.set("descr", "first floor plan")
    pptx = tmp_path / "floor.pptx"
    presentation.save(pptx)

    parsed = ReferencePptxParser().parse(
        pptx, workspace_dir=tmp_path / "ws2", capture_screenshots=False
    )
    pics = [
        e
        for e in parsed.slides[0].elements
        if e.element_type in {ReferenceElementType.IMAGE, ReferenceElementType.DRAWING}
    ]
    assert pics[0].element_type == ReferenceElementType.DRAWING
