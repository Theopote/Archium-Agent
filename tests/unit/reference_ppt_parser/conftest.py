"""Shared helpers for building multi-page reference PPTX fixtures."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def write_architectural_reference_pptx(path: Path, *, pages: int = 16) -> Path:
    """Create a ≥15-page architectural-style reference deck for induction tests."""
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    blank = presentation.slide_layouts[6]

    specs: list[tuple[str, list[tuple[str, float, float, float, float, int]]]] = [
        (
            "cover",
            [("院区更新概念汇报", 0.8, 2.0, 8.0, 1.0, 32)],
        ),
        (
            "agenda",
            [
                ("目录", 0.7, 0.4, 4.0, 0.6, 28),
                ("01 现状问题\n02 空间策略\n03 实施路径", 0.7, 1.4, 6.0, 2.5, 16),
            ],
        ),
        (
            "section",
            [("一、现状问题", 1.5, 2.2, 7.0, 1.0, 30)],
        ),
        (
            "photo_analysis",
            [
                ("入口交通拥堵问题突出", 0.6, 0.35, 8.5, 0.6, 24),
                ("现场判断：车行与人行冲突", 0.6, 1.1, 4.5, 1.8, 14),
                ("问题证据 1", 5.5, 1.1, 3.8, 1.5, 12),
                ("问题证据 2", 5.5, 2.8, 3.8, 1.5, 12),
            ],
        ),
        (
            "photo_analysis",
            [
                ("后勤流线交叉影响运营", 0.6, 0.35, 8.5, 0.6, 24),
                ("现场观察：装卸与急诊冲突", 0.6, 1.1, 4.5, 1.8, 14),
                ("问题证据 A", 5.5, 1.1, 3.8, 1.5, 12),
                ("问题证据 B", 5.5, 2.8, 3.8, 1.5, 12),
            ],
        ),
        (
            "drawing_focus",
            [
                ("总平面结构与关键指标", 0.6, 0.35, 8.5, 0.6, 24),
                ("总平面图", 0.6, 1.2, 5.5, 3.8, 14),
                ("总建筑面积 12.5 万㎡", 6.4, 1.3, 3.0, 0.7, 14),
                ("绿地率 35%", 6.4, 2.2, 3.0, 0.7, 14),
                ("来源：设计文本", 6.4, 4.5, 3.0, 0.5, 10),
            ],
        ),
        (
            "drawing_focus",
            [
                ("首层平面与流线组织", 0.6, 0.35, 8.5, 0.6, 24),
                ("平面图", 0.6, 1.2, 5.5, 3.8, 14),
                ("门诊面积 2.1 万㎡", 6.4, 1.3, 3.0, 0.7, 14),
                ("床位数 800", 6.4, 2.2, 3.0, 0.7, 14),
            ],
        ),
        (
            "metric",
            [
                ("经济技术指标摘要", 0.6, 0.35, 8.5, 0.6, 24),
                ("用地面积 6.8 万㎡", 0.8, 1.4, 4.0, 0.8, 16),
                ("容积率 2.1", 5.2, 1.4, 4.0, 0.8, 16),
                ("停车位 1200", 0.8, 2.6, 4.0, 0.8, 16),
                ("绿地率 35%", 5.2, 2.6, 4.0, 0.8, 16),
            ],
        ),
        (
            "metric",
            [
                ("分期实施关键指标", 0.6, 0.35, 8.5, 0.6, 24),
                ("一期投资 8.5 亿", 0.8, 1.4, 4.0, 0.8, 16),
                ("二期投资 6.2 亿", 5.2, 1.4, 4.0, 0.8, 16),
                ("总工期 48 个月", 0.8, 2.6, 4.0, 0.8, 16),
            ],
        ),
        (
            "section",
            [("二、空间策略", 1.5, 2.2, 7.0, 1.0, 30)],
        ),
        (
            "strategy",
            [
                ("三条空间策略形成清晰结构", 0.6, 0.35, 8.5, 0.6, 24),
                ("策略一：疏解入口交通", 0.7, 1.3, 8.5, 0.8, 16),
                ("策略二：分离后勤与急诊", 0.7, 2.3, 8.5, 0.8, 16),
                ("策略三：强化景观轴线", 0.7, 3.3, 8.5, 0.8, 16),
            ],
        ),
        (
            "before_after",
            [
                ("入口空间前后对比", 0.6, 0.35, 8.5, 0.6, 24),
                ("改造前", 0.7, 1.2, 4.0, 3.2, 14),
                ("改造后", 5.2, 1.2, 4.0, 3.2, 14),
            ],
        ),
        (
            "case_comparison",
            [
                ("对标案例对比", 0.6, 0.35, 8.5, 0.6, 24),
                ("案例 A", 0.7, 1.2, 4.0, 3.2, 14),
                ("案例 B", 5.2, 1.2, 4.0, 3.2, 14),
            ],
        ),
        (
            "grid",
            [
                ("现场调研照片网格", 0.6, 0.35, 8.5, 0.6, 24),
                ("照片1", 0.5, 1.2, 2.1, 1.8, 11),
                ("照片2", 2.8, 1.2, 2.1, 1.8, 11),
                ("照片3", 5.1, 1.2, 2.1, 1.8, 11),
                ("照片4", 7.4, 1.2, 2.1, 1.8, 11),
            ],
        ),
        (
            "decision",
            [
                ("提请决策事项", 0.6, 0.35, 8.5, 0.6, 24),
                ("建议批准一期实施方案与投资上限", 0.7, 1.5, 8.5, 2.0, 18),
            ],
        ),
        (
            "closing",
            [("谢谢", 3.5, 2.2, 3.0, 1.0, 36)],
        ),
    ]

    # Repeat content patterns if more pages requested.
    while len(specs) < pages:
        specs.append(specs[3 + (len(specs) % 6)])

    for _kind, boxes in specs[:pages]:
        slide = presentation.slides.add_slide(blank)
        for text, x, y, w, h, size in boxes:
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            frame = box.text_frame
            frame.text = text
            if frame.paragraphs and frame.paragraphs[0].runs:
                frame.paragraphs[0].runs[0].font.size = Pt(size)
                frame.paragraphs[0].runs[0].font.name = "Microsoft YaHei"

    # Add an anomalous dense page that should not win as representative.
    if pages >= 12:
        slide = presentation.slides.add_slide(blank)
        for i in range(22):
            box = slide.shapes.add_textbox(
                Inches(0.3 + (i % 4) * 2.4),
                Inches(0.3 + (i // 4) * 0.9),
                Inches(2.2),
                Inches(0.7),
            )
            box.text_frame.text = f"异常碎片 {i}"

    presentation.save(path)
    return path
