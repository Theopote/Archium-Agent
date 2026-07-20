"""Unit/integration tests for Template Studio import → annotate → fill → publish."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from archium.application.visual.template_studio_service import TemplateStudioService
from archium.config.settings import Settings
from archium.domain.visual.architectural_template import (
    TemplatePageType,
    TemplateSlot,
    TemplateSlotRole,
    TemplateStatus,
)
from archium.domain.visual.enums import DesignSystemSource
from archium.infrastructure.database.visual_repositories import DesignSystemRepository


def _write_sample_pptx(path: Path) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.7))
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "模板测试封面"
    run.font.size = Pt(32)
    run.font.name = "Arial"
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.0), Inches(2.0))
    body.text_frame.text = "用于 Template Studio 验收的正文区域。"
    presentation.save(path)
    return path


def test_template_studio_import_annotate_fill_publish(
    db_session: Session,
    tmp_path: Path,
) -> None:
    settings = Settings(_env_file=None, output_path=tmp_path)
    pptx_path = _write_sample_pptx(tmp_path / "source.pptx")
    service = TemplateStudioService(db_session, settings=settings)

    imported = service.import_pptx(pptx_path, name="验收模板")
    template = imported.template
    assert template.status == TemplateStatus.DRAFT
    assert template.layouts
    assert template.design_system_id is not None
    design = DesignSystemRepository(db_session).get(template.design_system_id)
    assert design is not None
    assert design.source_type == DesignSystemSource.IMPORTED
    assert template.fonts

    layout = template.layouts[0]
    template = service.update_page_type(template.id, layout.id, TemplatePageType.COVER)
    layout = template.layout_by_id(layout.id)
    assert layout is not None
    assert layout.page_type == TemplatePageType.COVER

    new_slot = TemplateSlot(
        id="manual_title",
        role=TemplateSlotRole.TITLE,
        required=True,
        x=0.8,
        y=0.4,
        width=8.0,
        height=0.8,
        auto_detected=False,
        label="人工标题槽",
    )
    template = service.upsert_slot(template.id, layout.id, new_slot)
    layout = template.layout_by_id(layout.id)
    assert layout is not None
    assert any(slot.id == "manual_title" for slot in layout.slots)

    preview = service.fill_test_content_preview(template.id, layout.id)
    assert preview.preview_path.is_file()
    assert preview.layout_plan.elements

    published = service.publish(template.id)
    assert published.status == TemplateStatus.PUBLISHED
    reloaded = service.get_template(template.id)
    assert reloaded is not None
    assert reloaded.status == TemplateStatus.PUBLISHED
    assert (Path(reloaded.workspace_dir) / "template.json").is_file()
