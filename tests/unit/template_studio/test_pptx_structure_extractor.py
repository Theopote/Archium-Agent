"""Unit tests for PPTX structure extraction used by Template Studio."""

from __future__ import annotations

from pathlib import Path

from archium.domain.visual.architectural_template import TemplatePageType, TemplateSlotRole
from archium.infrastructure.template.pptx_structure_extractor import PptxStructureExtractor
from pptx import Presentation
from pptx.util import Inches, Pt


def _write_sample_pptx(path: Path) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "院区总平面与指标"
    title_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title_frame.paragraphs[0].runs[0].font.name = "Microsoft YaHei"
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(4.5), Inches(2.5))
    body.text_frame.text = "正文说明：空间策略与实施路径。"
    body.text_frame.paragraphs[0].runs[0].font.name = "Microsoft YaHei"
    metric = slide.shapes.add_textbox(Inches(6.0), Inches(1.5), Inches(3.0), Inches(1.0))
    metric.text_frame.text = "总建筑面积 12.5 万㎡"
    presentation.save(path)
    return path


def test_pptx_structure_extractor_reads_fonts_slots_and_classification(tmp_path: Path) -> None:
    pptx_path = _write_sample_pptx(tmp_path / "sample.pptx")
    extraction = PptxStructureExtractor().extract(pptx_path)
    assert extraction.metadata.slide_count == 1
    assert extraction.pages
    assert extraction.fonts
    assert "Microsoft YaHei" in extraction.fonts
    page = extraction.pages[0]
    assert page.slots
    roles = {slot.role for slot in page.slots}
    assert TemplateSlotRole.TITLE in roles or TemplateSlotRole.BODY in roles
    layouts = PptxStructureExtractor().to_layouts(extraction)
    assert layouts[0].page_type in {
        TemplatePageType.COVER,
        TemplatePageType.TEXT_ARGUMENT,
        TemplatePageType.METRIC,
        TemplatePageType.UNKNOWN,
        TemplatePageType.SECTION,
    }
