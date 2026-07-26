"""Seed a minimal presentation + PPTX for UI deliver rehearsal (no LLM)."""

from __future__ import annotations

from pathlib import Path

from archium.domain.enums import PresentationStatus, SlideStatus, SlideType
from archium.domain.presentation import Presentation
from archium.domain.slide import SlideSpec
from archium.infrastructure.database.repositories import (
    PresentationRepository,
    ProjectRepository,
)
from archium.infrastructure.database.session import get_session
from pptx import Presentation as PptxPresentation


def main() -> None:
    with get_session() as session:
        projects = ProjectRepository(session).list_all()
        if not projects:
            raise SystemExit("No projects in DB")
        project = next((p for p in projects if "黄土" in p.name), projects[0])
        print(f"project={project.id} name={project.name}")

        presentations = PresentationRepository(session).list_by_project(project.id)
        if presentations:
            presentation = presentations[0]
            print(f"reuse presentation={presentation.id} title={presentation.title}")
        else:
            presentation = PresentationRepository(session).create_presentation(
                Presentation(
                    project_id=project.id,
                    title=f"{project.name} · UI验收稿",
                    status=PresentationStatus.DRAFT,
                )
            )
            print(f"created presentation={presentation.id}")

        slides = PresentationRepository(session).list_slides(presentation.id)
        if not slides:
            for order, title, stype, message, points in (
                (
                    0,
                    "项目定位",
                    SlideType.TITLE,
                    "黄土高原小型文化中心概念探索",
                    ["以光与生土为线索", "探索地域性公共建筑的当代表达"],
                ),
                (
                    1,
                    "核心问题",
                    SlideType.CONTENT,
                    "在干旱强光气候下，如何用建筑塑造可停留的公共光环境？",
                    ["气候与材料", "社区日常使用", "访客体验"],
                ),
                (
                    2,
                    "概念方向：黄土棱镜",
                    SlideType.CONTENT,
                    "光作为建造材料",
                    ["几何棱镜体块捕捉并折射光线", "形成随时间变化的室内光影"],
                ),
            ):
                PresentationRepository(session).save_slide(
                    SlideSpec(
                        presentation_id=presentation.id,
                        chapter_id="ch-concept",
                        order=order,
                        title=title,
                        message=message,
                        slide_type=stype,
                        status=SlideStatus.DRAFT,
                        key_points=list(points),
                        logical_key=f"ch-concept-p{order}",
                    )
                )
            slides = PresentationRepository(session).list_slides(presentation.id)
            print(f"created slides={len(slides)}")
        else:
            print(f"existing slides={len(slides)}")
        session.commit()
        presentation_id = presentation.id
        titles = [s.title for s in slides]

    out_dir = Path("data/outputs/ui_rehearsal")
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / f"{presentation_id}_ui_rehearsal.pptx"

    prs = PptxPresentation()
    for i, title in enumerate(titles or ["项目定位", "核心问题", "概念方向"]):
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if i > 0 and len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Archium UI 验收种子页（规则模式，无 LLM）"
    prs.save(pptx_path)
    print(f"DONE presentation_id={presentation_id} pptx={pptx_path} size={pptx_path.stat().st_size}")


if __name__ == "__main__":
    main()
