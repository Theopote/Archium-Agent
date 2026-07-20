"""Materialize desensitized cultural-village drop-in files for Phase 7 acceptance.

All content uses the fictional case name 砚溪村（脱敏）. Replace individual files with
real sanitized project exports when available; re-run this script only to regenerate stubs.
"""

from __future__ import annotations

import argparse
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation

from archium.domain.enums import VisualType
from archium.infrastructure.renderers.diagram_generator import generate_fallback_diagram

VILLAGE_NAME = "砚溪村（脱敏）"
_REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_ROOT = _REPO_ROOT / "tests" / "e2e" / "real_projects" / "files" / "cultural_village_001"

ASSET_SPECS: list[tuple[str, str, str]] = [
    ("01_village_aerial.png", "村落航拍", "photo"),
    ("02_historic_alley.png", "历史街巷", "photo"),
    ("03_ancestral_hall_plan.png", "宗祠平面", "floor_plan"),
    ("04_water_network_plan.png", "水系格局", "site_plan"),
    ("05_courtyard_life.png", "院落生活场景", "photo"),
    ("06_heritage_elements.png", "文保要素分布", "diagram"),
    ("07_tourism_circulation.png", "旅游流线组织", "diagram"),
    ("08_repair_node_detail.png", "修缮节点示意", "diagram"),
    ("09_activation_strategy_plan.png", "活化策略平面", "floor_plan"),
    ("10_phasing_diagram.png", "分期实施示意", "timeline"),
    ("11_street_section.png", "街巷断面示意", "diagram"),
]

_VISUAL_TYPE_MAP: dict[str, VisualType] = {
    "site_plan": VisualType.SITE_PLAN,
    "floor_plan": VisualType.FLOOR_PLAN,
    "diagram": VisualType.DIAGRAM,
    "timeline": VisualType.TIMELINE,
}


def _write_docx(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    document.add_heading(f"{VILLAGE_NAME}村落调研纪要", level=1)
    sections = [
        (
            "村落概况",
            [
                f"{VILLAGE_NAME}位于江南水网平原，形成于明清时期，现存街巷格局与宗祠、码头遗存较为完整。",
                "本次调研以现场踏勘、村民访谈与地方志摘录为主，不涉及真实地名与个人隐私。",
            ],
        ),
        (
            "水系与街巷",
            [
                "村内主河道呈 Y 形串联祠堂前广场与集市节点，支渠渗透至院落腹地。",
                "旅游旺季主街高峰人流约为平日 3 倍，居民反映噪声与垃圾清运压力显著。",
            ],
        ),
        (
            "宗祠与公共建筑",
            [
                "宗祠三进院落保存较好，梁架局部糟朽，需结构检测后分级修缮。",
                "戏台与晒场仍承担节庆活动，但缺少明确的分区管理与导览标识。",
            ],
        ),
        (
            "民居与活化",
            [
                "约 38% 民居空置或仅季节性居住，结构老化与设施不足并存。",
                "建议以“保护底线 + 分层活化”为原则，优先修复沿水巷的连续界面。",
            ],
        ),
    ]
    for title, paragraphs in sections:
        document.add_heading(title, level=2)
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
    document.save(path)


def _write_pdf(path: Path, *, title: str, paragraphs: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    document = fitz.open()
    page = document.new_page()
    y = 72
    page.insert_text((72, y), title, fontsize=14)
    y += 28
    for paragraph in paragraphs:
        page.insert_text((72, y), paragraph, fontsize=11)
        y += 18
        if y > 720:
            page = document.new_page()
            y = 72
    document.save(path)
    document.close()


def _write_xlsx(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "村落基础指标"
    rows = [
        ["指标", "数值", "备注"],
        ["案例化名", VILLAGE_NAME, "脱敏"],
        ["文保等级", "省级历史文化名村", "摘录"],
        ["常住户数", "186 户", "2024 调研"],
        ["户籍人口", "512 人", "摘录"],
        ["建设用地", "42.6 ha", "示意"],
        ["传统街巷", "1.8 km", "连续界面"],
        ["文保建筑", "17 处", "含宗祠"],
        ["年游客量", "约 28 万人次", "旺季集中"],
        ["空置民居比例", "38%", "需分级活化"],
    ]
    for row in rows:
        sheet.append(row)
    workbook.save(path)


def _write_pptx(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_layout)
    slide.shapes.title.text = f"{VILLAGE_NAME}参考汇报版式"
    slide.placeholders[1].text = "脱敏参考 · 版式与章节节奏示意（非最终方案）"

    for section in ["村落背景", "文化叙事", "保护策略", "活化路径"]:
        body_layout = presentation.slide_layouts[1]
        content = presentation.slides.add_slide(body_layout)
        content.shapes.title.text = section
        content.placeholders[1].text = f"{section}版式参考：标题区 + 主图区 + 要点栏"
    presentation.save(path)


def _write_photo_asset(path: Path, *, title: str, tint: tuple[int, int, int], filename: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    width, height = 1280, 800
    image = Image.new("RGB", (width, height), color=tint)
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, height - 120, width, height), fill=(30, 30, 30))
    draw.text((48, height - 80), title, fill=(245, 245, 240))
    draw.text((48, height - 48), f"{VILLAGE_NAME} · {filename}", fill=(200, 200, 195))
    image.save(path, format="PNG")


def _write_diagram_asset(path: Path, *, title: str, kind: str, filename: str) -> None:
    visual_type = _VISUAL_TYPE_MAP.get(kind, VisualType.DIAGRAM)
    generate_fallback_diagram(
        path,
        title=title,
        visual_type=visual_type,
        description=f"{VILLAGE_NAME} · {filename}",
        key_points=[
            f"{filename} · {title} · 控制点 A",
            f"{filename} · {title} · 控制点 B",
            f"{filename} · {title} · 控制点 C",
        ],
        message=f"{filename}::{title}",
    )
    # Ensure perceptual-hash dedup does not collapse distinct fixture assets.
    image = Image.open(path)
    draw = ImageDraw.Draw(image)
    draw.text((24, image.height - 36), filename, fill=(40, 40, 40))
    image.save(path, format="PNG")


def materialize_package(root: Path) -> list[Path]:
    """Create all drop-in files; returns written paths."""
    written: list[Path] = []

    docx_path = root / "documents" / "村落调研纪要.docx"
    _write_docx(docx_path)
    written.append(docx_path)

    pdf_research = root / "documents" / "文化价值研究摘要.pdf"
    _write_pdf(
        pdf_research,
        title=f"{VILLAGE_NAME}文化价值研究摘要",
        paragraphs=[
            "村落价值体现在水系街巷格局、宗祠礼制空间与民居建造技艺的连续性。",
            "文化叙事应围绕“同源共居—礼序生活—当代转译”三层展开。",
            "对外展示需避免将村民生活场景过度景观化。",
        ],
    )
    written.append(pdf_research)

    pdf_heritage = root / "documents" / "文保划定说明.pdf"
    _write_pdf(
        pdf_heritage,
        title=f"{VILLAGE_NAME}文保划定说明（摘录）",
        paragraphs=[
            "核心保护范围：宗祠—前广场—主河道一线。",
            "建设控制地带：主街两侧 30 m 缓冲及码头节点。",
            "禁止大拆大建，立面整治须保留坡屋顶与粉墙黛瓦特征。",
        ],
    )
    written.append(pdf_heritage)

    pptx_path = root / "documents" / "参考汇报版式.pptx"
    _write_pptx(pptx_path)
    written.append(pptx_path)

    xlsx_path = root / "data" / "村落基础指标.xlsx"
    _write_xlsx(xlsx_path)
    written.append(xlsx_path)

    photo_tints = [
        (92, 118, 104),
        (118, 102, 88),
        (104, 112, 128),
    ]
    photo_index = 0
    for filename, label, kind in ASSET_SPECS:
        asset_path = root / "assets" / filename
        if kind == "photo":
            _write_photo_asset(
                asset_path,
                title=label,
                tint=photo_tints[photo_index % len(photo_tints)],
                filename=filename,
            )
            photo_index += 1
        else:
            _write_diagram_asset(asset_path, title=label, kind=kind, filename=filename)
        written.append(asset_path)

    return written


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--root",
        type=Path,
        default=DEFAULT_ROOT,
        help="Output directory (default: tests/e2e/real_projects/files/cultural_village_001)",
    )
    args = parser.parse_args()
    paths = materialize_package(args.root.resolve())
    print(f"Wrote {len(paths)} files under {args.root.resolve()}")


if __name__ == "__main__":
    main()
