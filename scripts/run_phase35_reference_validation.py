"""Phase 3.5: run real reference PPTX induction and emit a human review checklist.

Usage (from repo root):
    py scripts/run_phase35_reference_validation.py
    py scripts/run_phase35_reference_validation.py --source "C:\\path\\to\\deck.pptx"
    py scripts/run_phase35_reference_validation.py --max-slides 28 --require-screenshots
"""

from __future__ import annotations

import argparse
import json
import shutil
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation

# Repo root on sys.path
_REPO = Path(__file__).resolve().parents[1]
if str(_REPO) not in sys.path:
    sys.path.insert(0, str(_REPO))

from archium.application.visual.asset_path_resolver import is_machine_absolute_path
from archium.application.visual.template_induction_service import TemplateInductionService
from archium.domain.visual.reference_slide import ReferenceElementType
from archium.domain.visual.template_induction import FunctionalSlideType
from archium.infrastructure.renderers.pptx_screenshot import screenshot_tools_available


DEFAULT_SOURCE_DIR = Path(r"C:\Users\navib\Desktop\development\参考pptx")


def _pick_source_pptx(source_dir: Path, explicit: Path | None) -> Path:
    if explicit is not None:
        if not explicit.is_file():
            raise FileNotFoundError(explicit)
        return explicit
    candidates = sorted(source_dir.glob("*.pptx"))
    if not candidates:
        raise FileNotFoundError(f"No .pptx in {source_dir}")
    # Prefer smallest real deck in the folder (still architectural).
    return min(candidates, key=lambda p: len(Presentation(str(p)).slides))


def _subset_pptx(source: Path, dest: Path, *, max_slides: int) -> tuple[Path, int]:
    prs = Presentation(str(source))
    total = len(prs.slides)
    if total <= max_slides:
        shutil.copy2(source, dest)
        return dest, total
    out = Presentation(str(source))
    # python-pptx: drop from end until count matches.
    while len(out.slides) > max_slides:
        r_id = out.slides._sldIdLst[-1].rId  # noqa: SLF001
        out.part.drop_rel(r_id)
        del out.slides._sldIdLst[-1]
    out.save(dest)
    return dest, max_slides


def _fingerprint_clusters(induction) -> list[tuple]:
    return sorted(
        (
            c.functional_type.value,
            c.content_type.value,
            tuple(c.slide_ids),
            c.representative_slide_id,
        )
        for c in induction.clusters
        if c.functional_type == FunctionalSlideType.CONTENT
    )


def run_validation(
    *,
    source_dir: Path,
    explicit_source: Path | None,
    max_slides: int,
    require_screenshots: bool,
    output_root: Path,
) -> dict:
    source = _pick_source_pptx(source_dir, explicit_source)
    source_slides = len(Presentation(str(source)).slides)

    stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    run_dir = output_root / f"phase35_{stamp}"
    run_dir.mkdir(parents=True, exist_ok=True)

    working = run_dir / "input.pptx"
    working, used_slides = _subset_pptx(source, working, max_slides=max_slides)

    tools = screenshot_tools_available()
    if require_screenshots and not tools:
        raise RuntimeError("Screenshot tools unavailable but --require-screenshots set")

    service = TemplateInductionService()
    service.workspace_root = lambda induction_id: run_dir / "induction" / str(induction_id)  # type: ignore[method-assign]

    first = service.induce(
        working,
        name=working.stem,
        capture_screenshots=True,
        require_screenshots=require_screenshots,
    )
    second = service.induce(
        working,
        name=working.stem + "_rerun",
        capture_screenshots=False,
        require_screenshots=False,
    )

    presentation = first.presentation
    induction = first.induction
    workspace = first.workspace

    content_clusters = [
        c for c in induction.clusters if c.functional_type == FunctionalSlideType.CONTENT
    ]
    types = Counter(c.functional_type for c in induction.classifications)
    drawing_elems = sum(
        1
        for slide in presentation.slides
        for elem in slide.iter_elements()
        if elem.element_type == ReferenceElementType.DRAWING
    )

    png_missing = [
        s.slide_id
        for s in presentation.slides
        if not s.image_path or not (workspace / s.image_path).is_file()
    ]
    abs_paths: list[str] = []
    payload = json.loads((workspace / "reference_presentation.json").read_text(encoding="utf-8"))
    for slide in payload.get("slides", []):
        for key in ("image_path",):
            val = slide.get(key) or ""
            if is_machine_absolute_path(val):
                abs_paths.append(val)
        for asset in slide.get("image_assets") or []:
            rel = asset.get("relative_path") or ""
            if is_machine_absolute_path(rel):
                abs_paths.append(rel)

    rep_scores = {s.slide_id: s for s in induction.representative_scores}
    anomalous_reps = []
    for cluster in induction.clusters:
        rep = cluster.representative_slide_id
        if not rep:
            continue
        score = rep_scores.get(rep)
        if score and (score.anomaly_penalty > 0.2 or score.excessive_complexity_penalty > 0.2):
            anomalous_reps.append(
                {
                    "cluster_id": cluster.id,
                    "representative_slide_id": rep,
                    "anomaly_penalty": score.anomaly_penalty,
                    "complexity_penalty": score.excessive_complexity_penalty,
                }
            )

    stable = _fingerprint_clusters(first.induction) == _fingerprint_clusters(second.induction)

    checks = {
        "slide_count_match": presentation.slide_count == used_slides == induction.slide_count,
        "screenshot_count_match": first.screenshot_count == presentation.slide_count,
        "all_png_present": len(png_missing) == 0,
        "functional_cover": FunctionalSlideType.COVER in types,
        "functional_content": FunctionalSlideType.CONTENT in types,
        "functional_closing_or_appendix": types.get(FunctionalSlideType.CLOSING, 0)
        + types.get(FunctionalSlideType.APPENDIX, 0)
        > 0,
        "content_clusters_gte_3": len(content_clusters) >= 3,
        "all_clusters_have_representative": all(c.representative_slide_id for c in induction.clusters),
        "drawing_elements_found": drawing_elems > 0,
        "no_machine_absolute_paths": len(abs_paths) == 0,
        "reference_assets_origin": all(
            asset.asset_origin == "reference_template"
            for slide in presentation.slides
            for asset in slide.image_assets
        ),
        "rerun_cluster_stable": stable,
        "anomalous_representative_count": len(anomalous_reps),
    }
    auto_pass = all(
        checks[k]
        for k in checks
        if k not in {"anomalous_representative_count", "functional_closing_or_appendix"}
    ) and checks["anomalous_representative_count"] == 0

    report = {
        "generated_at_utc": datetime.now(timezone.utc).isoformat(),
        "source_file": str(source),
        "source_slide_count": source_slides,
        "validated_slide_count": used_slides,
        "working_pptx": str(working),
        "workspace": str(workspace),
        "screenshot_tools_available": tools,
        "screenshot_count": first.screenshot_count,
        "functional_type_counts": {k.value: v for k, v in types.items()},
        "content_cluster_count": len(content_clusters),
        "drawing_element_count": drawing_elems,
        "low_confidence_slides": induction.low_confidence_slide_ids,
        "checks": checks,
        "auto_checks_pass": auto_pass,
        "anomalous_representatives": anomalous_reps,
        "content_clusters": [
            {
                "cluster_id": c.id,
                "content_type": c.content_type.value,
                "size": len(c.slide_ids),
                "representative_slide_id": c.representative_slide_id,
                "slide_ids": c.slide_ids,
                "rationale": c.selection_rationale[:6],
            }
            for c in content_clusters
        ],
        "human_review_checklist": _human_checklist(),
    }

    report_path = run_dir / "phase35_validation_report.json"
    checklist_path = run_dir / "phase35_human_review_checklist.md"
    report_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2, default=str),
        encoding="utf-8",
    )
    checklist_path.write_text(_render_checklist_md(report), encoding="utf-8")
    return report


def _human_checklist() -> list[dict[str, str]]:
    return [
        {
            "id": "H1",
            "item": "页数一致性",
            "question": "源 PPT 页数、解析页数、slide JSON 数、PNG 截图数是否一致？",
            "auto": "slide_count_match + screenshot_count_match + all_png_present",
        },
        {
            "id": "H2",
            "item": "功能页类型",
            "question": "封面/目录或章节/内容/结尾（或附录）是否基本合理？误把免责声明当封面？",
            "auto": "functional_* flags（需目视确认）",
        },
        {
            "id": "H3",
            "item": "内容聚类",
            "question": "至少 3 个内容聚类是否有语义意义？不应仅靠 layout_name 硬拆？",
            "auto": "content_clusters_gte_3",
        },
        {
            "id": "H4",
            "item": "图纸识别",
            "question": "总平面/平面图/剖面等是否标为 Drawing（非普通 Photo）？",
            "auto": "drawing_elements_found",
        },
        {
            "id": "H5",
            "item": "代表页",
            "question": "代表页是否结构完整、可编辑？异常密页/解析失败页不应成为代表页。",
            "auto": "anomalous_representative_count == 0",
        },
        {
            "id": "H6",
            "item": "参考隔离",
            "question": "参考页文字/图片是否未进入项目 Manuscript（仅 reference_template）？",
            "auto": "reference_assets_origin",
        },
        {
            "id": "H7",
            "item": "路径可移植",
            "question": "产物 JSON 是否无本机绝对路径？",
            "auto": "no_machine_absolute_paths",
        },
        {
            "id": "H8",
            "item": "聚类修正",
            "question": "Review UI 能否移动/合并/拆分聚类并保存？",
            "auto": "manual — 在 Streamlit 模板归纳页试一次",
        },
        {
            "id": "H9",
            "item": "重跑稳定",
            "question": "同一 PPT 两次归纳，内容聚类成员与代表页是否一致？",
            "auto": "rerun_cluster_stable",
        },
    ]


def _render_checklist_md(report: dict) -> str:
    lines = [
        "# Phase 3.5 人工复核清单",
        "",
        f"- 生成时间（UTC）：{report['generated_at_utc']}",
        f"- 源文件：`{report['source_file']}`",
        f"- 源页数：{report['source_slide_count']}",
        f"- 本次验收页数：{report['validated_slide_count']}",
        f"- 工作区：`{report['workspace']}`",
        f"- 自动检查通过：{'是' if report['auto_checks_pass'] else '否'}",
        "",
        "## 自动检查结果",
        "",
        "| 检查项 | 结果 |",
        "|--------|------|",
    ]
    for key, value in report["checks"].items():
        lines.append(f"| `{key}` | {value} |")
    lines.extend(
        [
            "",
            "## 内容聚类摘要",
            "",
        ]
    )
    for cluster in report["content_clusters"]:
        lines.append(
            f"- **{cluster['cluster_id'][:8]}** · {cluster['content_type']} · "
            f"{cluster['size']} 页 · 代表 `{cluster['representative_slide_id']}`"
        )
    if report["low_confidence_slides"]:
        lines.extend(["", "## 待复核页", ""])
        for sid in report["low_confidence_slides"]:
            lines.append(f"- `{sid}`")
    lines.extend(["", _LOW_VISUAL_DENSITY_RUBRIC_MD, "", "## 人工逐项确认", "", "| ID | 项目 | 自动 | 人工结论 | 备注 |", "|----|------|------|----------|------|"])
    for item in report["human_review_checklist"]:
        lines.append(
            f"| {item['id']} | {item['item']} | {item['auto']} | ☐ PASS / ☐ REVIEW | |"
        )
    lines.extend(
        [
            "",
            "## 签收",
            "",
            "- 复核人：",
            "- 日期：",
            "- 总体结论：☐ PASS  ☐ PASS_WITH_WARNINGS  ☐ NEEDS_REVIEW  ☐ BLOCKED",
            "",
        ]
    )
    return "\n".join(lines)


_LOW_VISUAL_DENSITY_RUBRIC_MD = """## 低视觉密度页评审口径（本批次必用）

适用对象：以展位/说明文字为主，几乎无建筑图纸或现场图，字体/颜色变化很少的页面。

评审原则：

- 不因“没有建筑图”自动扣分；先判断该页是否本来就是说明/目录/流程页。
- 弱化主观美学项（颜色丰富度、字体层次），改看信息结构与可读性。
- 优先判定 `PASS` / `REVIEW`，避免在该类页上做细粒度主观分差。

人工快速判定（建议）：

| 维度 | PASS 条件 | REVIEW 触发 |
|------|-----------|-------------|
| 结构清晰 | 标题、正文、强调层级可辨识 | 标题/正文混在一起，读者难以抓主线 |
| 可读性 | 正文密度可接受，行距和留白基本可读 | 文字过密、拥挤、明显难读 |
| 逻辑顺序 | 叙述顺序连贯（背景→问题→措施等） | 叙述跳跃、前后关系不清 |
| 编辑友好 | 文本槽位可替换，非大面积锁定/背景覆盖 | 主要信息位难改或需重做版面 |
| 与页面功能匹配 | 作为说明页/目录页目标明确 | 页功能与内容不匹配 |"""


def main() -> int:
    parser = argparse.ArgumentParser(description="Phase 3.5 real reference deck validation")
    parser.add_argument("--source-dir", type=Path, default=DEFAULT_SOURCE_DIR)
    parser.add_argument("--source", type=Path, default=None, help="Explicit .pptx path")
    parser.add_argument("--max-slides", type=int, default=28, help="Cap slides for sprint (15–30)")
    parser.add_argument("--require-screenshots", action="store_true", default=True)
    parser.add_argument("--no-screenshots", action="store_true")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=_REPO / "output" / "phase35-validation",
    )
    args = parser.parse_args()
    require_shots = args.require_screenshots and not args.no_screenshots

    report = run_validation(
        source_dir=args.source_dir,
        explicit_source=args.source,
        max_slides=args.max_slides,
        require_screenshots=require_shots,
        output_root=args.output_root,
    )
    print(json.dumps(
        {
            "auto_checks_pass": report["auto_checks_pass"],
            "workspace": report["workspace"],
            "validated_slide_count": report["validated_slide_count"],
            "content_clusters": report["content_cluster_count"],
            "drawing_elements": report["drawing_element_count"],
        },
        ensure_ascii=False,
        indent=2,
    ))
    return 0 if report["auto_checks_pass"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
