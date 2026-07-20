"""Materialize desensitized renovation drop-in files for Phase 7 acceptance.

All content uses the fictional case name 东华文创园（脱敏）. Replace individual files with
real sanitized project exports when available; re-run this script only to regenerate stubs.
"""

from __future__ import annotations

import argparse
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation

from archium.domain.enums import VisualType
from archium.infrastructure.renderers.diagram_generator import generate_fallback_diagram

SITE_NAME = "东华文创园（脱敏）"
_REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_ROOT = _REPO_ROOT / "tests" / "e2e" / "real_projects" / "files" / "renovation_001"

ASSET_SPECS: list[tuple[str, str, str]] = [
    ("01_factory_aerial.png", "厂区航拍", "photo"),
    ("02_facade_before.png", "改造前面貌", "photo"),
    ("03_structure_existing.png", "结构现状", "photo"),
    ("04_floor_plan_before.png", "改造前平面", "floor_plan"),
    ("05_floor_plan_after.png", "改造后平面", "floor_plan"),
    ("06_public_realm.png", "公共界面", "site_plan"),
    ("07_circulation_diagram.png", "人流组织", "diagram"),
    ("08_phasing_diagram.png", "分期示意", "timeline"),
    ("09_section_renovation.png", "改造剖面", "diagram"),
    ("10_ground_activation.png", "首层活化", "floor_plan"),
    ("11_roof_intervention.png", "屋顶加建示意", "diagram"),
]

_VISUAL_TYPE_MAP: dict[str, VisualType] = {
    "site_plan": VisualType.SITE_PLAN,
    "floor_plan": VisualType.FLOOR_PLAN,
    "diagram": VisualType.DIAGRAM,
    "timeline": VisualType.TIMELINE,
}


def _write_docx(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    document.add_heading(f"{SITE_NAME}改造任务书（脱敏摘录）", level=1)
    sections = [
        (
            "项目背景",
            [
                f"{SITE_NAME}原为 1980 年代工业厂房，现拟转型为创意办公与展示复合功能。",
                "业主希望保留主体结构框架，提升公共界面与首层开放性，分三期实施。",
            ],
        ),
        (
            "现状问题",
            [
                "屋面渗漏与局部构件锈蚀需优先处理；消防疏散不满足现行规范。",
                "厂区主入口形象陈旧，缺乏连续公共界面；停车与卸货流线交叉。",
            ],
        ),
        (
            "改造原则",
            [
                "结构能留则留，局部加固与置换并重；机电系统整体更新。",
                "首层开放、上层弹性办公；保留工业遗存记忆点作为叙事锚点。",
            ],
        ),
        (
            "分期要求",
            [
                "一期保障既有租户运营，完成屋面与消防改造；",
                "二期完成公共界面与主入口；三期完善屋顶公共层与景观。",
            ],
        ),
    ]
    for title, paragraphs in sections:
        document.add_heading(title, level=2)
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
    document.save(path)


def _write_pdf(path: Path, *, title: str, paragraphs: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    document = fitz.open()
    page = document.new_page()
    y = 72
    page.insert_text((72, y), title, fontsize=14)
    y += 28
    for paragraph in paragraphs:
        page.insert_text((72, y), paragraph, fontsize=11)
        y += 18
        if y > 720:
            page = document.new_page()
            y = 72
    document.save(path)
    document.close()


def _write_xlsx(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "改造面积指标"
    rows = [
        ["指标", "现状", "改造后", "备注"],
        ["案例化名", SITE_NAME, SITE_NAME, "脱敏"],
        ["建筑面积", "2.4 万㎡", "3.1 万㎡", "含屋顶加建"],
        ["公共空间占比", "8%", "18%", "首层与屋顶"],
        ["办公可租面积", "1.6 万㎡", "2.2 万㎡", "弹性隔断"],
        ["停车位", "86 个", "120 个", "地下补充"],
        ["结构保留率", "—", "约 72%", "框架保留"],
        ["改造分期", "—", "三期", "24 个月"],
    ]
    for row in rows:
        sheet.append(row)
    workbook.save(path)


def _write_pptx(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_layout)
    slide.shapes.title.text = f"{SITE_NAME}改造汇报参考版式"
    slide.placeholders[1].text = "脱敏参考 · 版式与章节节奏示意（非最终方案）"

    for section in ["厂区现状", "问题诊断", "改造策略", "分期实施"]:
        body_layout = presentation.slide_layouts[1]
        content = presentation.slides.add_slide(body_layout)
        content.shapes.title.text = section
        content.placeholders[1].text = f"{section}版式参考：对比图 + 策略要点 + 指标栏"
    presentation.save(path)


def _write_photo_asset(path: Path, *, title: str, tint: tuple[int, int, int], filename: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    width, height = 1280, 800
    image = Image.new("RGB", (width, height), color=tint)
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, height - 120, width, height), fill=(30, 30, 30))
    draw.text((48, height - 80), title, fill=(245, 245, 240))
    draw.text((48, height - 48), f"{SITE_NAME} · {filename}", fill=(200, 200, 195))
    image.save(path, format="PNG")


def _write_diagram_asset(path: Path, *, title: str, kind: str, filename: str) -> None:
    visual_type = _VISUAL_TYPE_MAP.get(kind, VisualType.DIAGRAM)
    generate_fallback_diagram(
        path,
        title=title,
        visual_type=visual_type,
        description=f"{SITE_NAME} · {filename}",
        key_points=[
            f"{filename} · {title} · 控制点 A",
            f"{filename} · {title} · 控制点 B",
            f"{filename} · {title} · 控制点 C",
        ],
        message=f"{filename}::{title}",
    )
    image = Image.open(path)
    draw = ImageDraw.Draw(image)
    draw.text((24, image.height - 36), filename, fill=(40, 40, 40))
    image.save(path, format="PNG")


def materialize_package(root: Path) -> list[Path]:
    written: list[Path] = []

    docx_path = root / "documents" / "改造任务书.docx"
    _write_docx(docx_path)
    written.append(docx_path)

    survey_pdf = root / "documents" / "现状调研报告.pdf"
    _write_pdf(
        survey_pdf,
        title=f"{SITE_NAME}现状调研报告（摘录）",
        paragraphs=[
            "厂区由三栋并联厂房组成，柱网 8.4 m × 8.4 m，层高首层 6 m、标准层 4.5 m。",
            "外立面为清水砖与金属窗框组合，局部已加装彩钢板围护。",
            "租户以设计工作室与小型制造为主，货运高峰集中在工作日上午。",
        ],
    )
    written.append(survey_pdf)

    structure_pdf = root / "documents" / "结构检测摘要.pdf"
    _write_pdf(
        structure_pdf,
        title=f"{SITE_NAME}结构检测摘要",
        paragraphs=[
            "框架梁柱总体完好，二层部分梁端锈蚀等级 II 级，建议除锈防腐并局部补强。",
            "屋面檩条老化，需整体更换；基础未见明显不均匀沉降。",
            "加建屋顶公共层需复核柱顶承载力，建议增设钢支撑桁架。",
        ],
    )
    written.append(structure_pdf)

    pptx_path = root / "documents" / "参考汇报版式.pptx"
    _write_pptx(pptx_path)
    written.append(pptx_path)

    xlsx_path = root / "data" / "改造面积指标.xlsx"
    _write_xlsx(xlsx_path)
    written.append(xlsx_path)

    photo_tints = [
        (88, 96, 108),
        (108, 96, 88),
        (96, 104, 112),
    ]
    photo_index = 0
    for filename, label, kind in ASSET_SPECS:
        asset_path = root / "assets" / filename
        if kind == "photo":
            _write_photo_asset(
                asset_path,
                title=label,
                tint=photo_tints[photo_index % len(photo_tints)],
                filename=filename,
            )
            photo_index += 1
        else:
            _write_diagram_asset(asset_path, title=label, kind=kind, filename=filename)
        written.append(asset_path)

    return written


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--root",
        type=Path,
        default=DEFAULT_ROOT,
        help="Output directory (default: tests/e2e/real_projects/files/renovation_001)",
    )
    args = parser.parse_args()
    paths = materialize_package(args.root.resolve())
    print(f"Wrote {len(paths)} files under {args.root.resolve()}")


if __name__ == "__main__":
    main()
