"""PPTX document parser."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide

from archium.domain.enums import AssetType
from archium.infrastructure.document_parsers._utils import (
    _PPTX_SUFFIXES,
    build_parsed_document,
    normalize_whitespace,
    suffix_of,
)
from archium.infrastructure.document_parsers.base import ExtractedAsset, ParsedDocument, ParsedPage


class PptxParser:
    """Extract slide text, notes, and images from PPTX files."""

    def supports(self, file_path: Path) -> bool:
        return suffix_of(file_path) in _PPTX_SUFFIXES

    def parse(self, file_path: Path) -> ParsedDocument:
        presentation = Presentation(str(file_path))
        pages: list[ParsedPage] = []
        assets: list[ExtractedAsset] = []

        slide_width = int(presentation.slide_width or 0)
        slide_height = int(presentation.slide_height or 0)
        metadata: dict[str, object] = {
            "slide_count": len(presentation.slides),
            "slide_width": slide_width,
            "slide_height": slide_height,
        }

        for slide_index, slide in enumerate(presentation.slides, start=1):
            title = self._slide_title(slide)
            body_lines: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = normalize_whitespace(paragraph.text)
                        if text:
                            body_lines.append(text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        assets.append(
                            ExtractedAsset(
                                filename=f"slide{slide_index}_{image.filename or 'image.png'}",
                                data=image.blob,
                                page_number=slide_index,
                                asset_type=AssetType.IMAGE,
                                description=f"Slide {slide_index} image",
                            )
                        )
                    except Exception:
                        continue

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = normalize_whitespace(slide.notes_slide.notes_text_frame.text)

            slide_text = normalize_whitespace("\n".join(body_lines))
            if notes:
                slide_text = normalize_whitespace(f"{slide_text}\n\n[Notes]\n{notes}")

            pages.append(
                ParsedPage(
                    page_number=slide_index,
                    text=slide_text or f"[Slide {slide_index} has no extractable text]",
                    section_title=title,
                    content_type="slide",
                )
            )

        return build_parsed_document(pages, metadata=metadata).model_copy(update={"assets": assets})

    def _slide_title(self, slide: Slide) -> str | None:
        if slide.shapes.title and slide.shapes.title.text:
            return normalize_whitespace(slide.shapes.title.text)
        return None
