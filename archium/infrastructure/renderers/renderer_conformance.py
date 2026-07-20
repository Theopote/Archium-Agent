"""Cross-renderer conformance checks for RenderScene outputs."""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from archium.domain.visual.render_scene import DrawingNode, ImageNode, RenderScene, TextNode
from archium.infrastructure.renderers.html_renderer import HtmlRenderer


@dataclass(frozen=True)
class RendererSnapshot:
    text_values: tuple[str, ...]
    image_node_ids: tuple[str, ...]
    node_count: int
    background_color: str


@dataclass
class ConformanceReport:
    passed: bool = True
    issues: list[str] = field(default_factory=list)

    def fail(self, message: str) -> None:
        self.passed = False
        self.issues.append(message)


def snapshot_from_scene(scene: RenderScene) -> RendererSnapshot:
    texts = tuple(
        node.text.strip()
        for node in scene.nodes
        if isinstance(node, TextNode) and node.text.strip()
    )
    images = tuple(
        node.id
        for node in scene.nodes
        if isinstance(node, (ImageNode, DrawingNode)) and not node.asset_unresolved
    )
    return RendererSnapshot(
        text_values=texts,
        image_node_ids=images,
        node_count=len(scene.nodes),
        background_color=scene.background.color.lower(),
    )


def snapshot_from_html(scene: RenderScene) -> RendererSnapshot:
    html = HtmlRenderer().render(scene)
    texts = tuple(
        match.strip()
        for match in re.findall(r'class="node text-node"[^>]*>([^<]+)<', html)
        if match.strip()
    )
    image_ids = tuple(re.findall(r'class="node image-node[^"]*" id="([^"]+)"', html))
    bg_match = re.search(r"background:\s*(#[0-9A-Fa-f]{6})", html)
    background = bg_match.group(1).lower() if bg_match else scene.background.color.lower()
    return RendererSnapshot(
        text_values=texts,
        image_node_ids=image_ids,
        node_count=len(re.findall(r'class="node', html)),
        background_color=background,
    )


def snapshot_from_pptx(path: Path) -> RendererSnapshot:
    from pptx import Presentation

    presentation = Presentation(path)
    texts: list[str] = []
    image_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "text", "") and str(shape.text).strip():
                texts.append(str(shape.text).strip())
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
    return RendererSnapshot(
        text_values=tuple(texts),
        image_node_ids=tuple(str(index) for index in range(image_count)),
        node_count=len(texts) + image_count,
        background_color="",
    )


def compare_snapshots(
    reference: RendererSnapshot,
    candidate: RendererSnapshot,
    *,
    label: str,
    position_tolerance_in: float = 0.08,
) -> ConformanceReport:
    del position_tolerance_in  # reserved for future geometry checks
    report = ConformanceReport()
    if set(reference.text_values) != set(candidate.text_values):
        report.fail(
            f"{label}: text mismatch "
            f"(ref={list(reference.text_values)!r}, got={list(candidate.text_values)!r})"
        )
    if len(reference.image_node_ids) != len(candidate.image_node_ids):
        report.fail(
            f"{label}: image count mismatch "
            f"(ref={len(reference.image_node_ids)}, got={len(candidate.image_node_ids)})"
        )
    return report


def assert_renderer_conformance(
    scene: RenderScene,
    *,
    pptx_path: Path | None = None,
) -> ConformanceReport:
    """Compare scene against HTML and optional PPTX exports."""
    reference = snapshot_from_scene(scene)
    report = ConformanceReport()
    html_report = compare_snapshots(reference, snapshot_from_html(scene), label="html")
    report.passed = report.passed and html_report.passed
    report.issues.extend(html_report.issues)

    if pptx_path is not None and pptx_path.is_file():
        pptx_report = compare_snapshots(
            reference,
            snapshot_from_pptx(pptx_path),
            label="pptx",
        )
        report.passed = report.passed and pptx_report.passed
        report.issues.extend(pptx_report.issues)
    return report
