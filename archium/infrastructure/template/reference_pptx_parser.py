"""Parse reference PPTX into ReferenceSlideSnapshot (induction IR).

Extends Template Studio extraction with richer element typing, signatures,
extracted reference assets, group recursion, and portable relative paths.
Does not invent a parallel generation kernel.
"""

from __future__ import annotations

import contextlib
import hashlib
import json
import re
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from archium.application.visual.drawing_inference_service import DrawingInferenceService
from archium.domain.visual.placeholder_binding import build_placeholder_binding_signature
from archium.domain.visual.reference_slide import (
    REFERENCE_TEMPLATE_ASSET_ORIGIN,
    ReferenceAsset,
    ReferenceElement,
    ReferenceElementType,
    ReferencePresentation,
    ReferenceSlideSnapshot,
)
from archium.infrastructure.renderers.pptx_screenshot import (
    export_pptx_slide_pngs,
    screenshot_tools_available,
)

_EMU_PER_INCH = 914400.0
_WINDOWS_ABS = re.compile(r"^[A-Za-z]:[\\/]")
_GENERIC_SHAPE_NAME = re.compile(
    r"^(textbox|rectangle|picture|oval|shape|group|placeholder|autoshape)\s*\d*$",
    re.I,
)
# Large non-text fills covering ~17.5%+ of the page area → decoration candidate.
_DECORATION_AREA_RATIO = 0.175


def _emu_to_inches(value: int | float | None) -> float:
    if value is None:
        return 0.0
    return round(float(value) / _EMU_PER_INCH, 4)


def _rgb_to_hex(rgb: object) -> str | None:
    try:
        return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"  # type: ignore[index]
    except Exception:
        return None


def _slide_id(index: int) -> str:
    return f"slide_{index + 1:03d}"


def _relative_slide_image(index: int) -> str:
    return f"slides/{_slide_id(index)}.png"


def _assert_relative(path: str) -> str:
    text = (path or "").strip().replace("\\", "/")
    if not text:
        return ""
    if text.startswith(("/", "\\\\")) or _WINDOWS_ABS.match(text):
        raise ValueError(f"must not persist machine absolute path: {path}")
    return text


def _content_signature(
    *,
    element_types: list[str],
    image_count: int,
    text_length: int,
    chart_count: int,
    table_count: int,
) -> str:
    """Structural signature — layout_name intentionally excluded (soft feature only)."""
    payload = {
        "types": element_types,
        "images": image_count,
        "text_len_bucket": text_length // 40,
        "charts": chart_count,
        "tables": table_count,
    }
    raw = json.dumps(payload, sort_keys=True, ensure_ascii=False)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:24]


def _axis_aligned_union_area(
    rects: list[tuple[float, float, float, float]],
    *,
    page_width: float,
    page_height: float,
) -> float:
    """Exact union area of axis-aligned rectangles clipped to the page.

    ``rects`` are ``(x, y, width, height)``. Overlaps are counted once.
    """
    clipped: list[tuple[float, float, float, float]] = []
    for x, y, w, h in rects:
        x0 = max(0.0, min(float(x), page_width))
        y0 = max(0.0, min(float(y), page_height))
        x1 = max(0.0, min(float(x) + float(w), page_width))
        y1 = max(0.0, min(float(y) + float(h), page_height))
        if x1 > x0 and y1 > y0:
            clipped.append((x0, y0, x1, y1))
    if not clipped:
        return 0.0
    xs = sorted({r[0] for r in clipped} | {r[2] for r in clipped})
    ys = sorted({r[1] for r in clipped} | {r[3] for r in clipped})
    area = 0.0
    for i in range(len(xs) - 1):
        for j in range(len(ys) - 1):
            cx0, cx1 = xs[i], xs[i + 1]
            cy0, cy1 = ys[j], ys[j + 1]
            if cx1 <= cx0 or cy1 <= cy0:
                continue
            mx = (cx0 + cx1) * 0.5
            my = (cy0 + cy1) * 0.5
            if any(r[0] <= mx <= r[2] and r[1] <= my <= r[3] for r in clipped):
                area += (cx1 - cx0) * (cy1 - cy0)
    return area


def _visual_embedding(
    *,
    width: float,
    height: float,
    elements: list[ReferenceElement],
    image_count: int,
    text_length: int,
    chart_count: int,
    table_count: int,
) -> list[float]:
    """Deterministic structural embedding for clustering (not a neural model)."""
    area = max(width * height, 0.01)
    flat = [node for e in elements for node in e.iter_self_and_descendants()]
    type_counts = Counter(e.element_type.value for e in flat)
    rects = [(e.x, e.y, e.width, e.height) for e in flat]
    # True page coverage (union), not summed element areas.
    covered = _axis_aligned_union_area(rects, page_width=width, page_height=height) / area
    top_band = height * 0.33
    left_band = width * 0.45
    top_rects = [
        (x, y, w, min(y + h, top_band) - y)
        for x, y, w, h in rects
        if y < top_band and min(y + h, top_band) > y
    ]
    left_rects = [
        (x, y, min(x + w, left_band) - x, h)
        for x, y, w, h in rects
        if x < left_band and min(x + w, left_band) > x
    ]
    top_heavy = (
        _axis_aligned_union_area(top_rects, page_width=width, page_height=top_band) / area
    )
    left_heavy = (
        _axis_aligned_union_area(left_rects, page_width=left_band, page_height=height)
        / area
    )
    max_element = max((e.width * e.height for e in flat), default=0.0) / area
    return [
        round(min(covered, 1.0), 4),
        round(min(top_heavy, 1.0), 4),
        round(min(left_heavy, 1.0), 4),
        round(min(max_element, 1.5), 4),
        round(image_count / 8.0, 4),
        round(min(text_length, 800) / 800.0, 4),
        round(chart_count / 3.0, 4),
        round(table_count / 2.0, 4),
        round(type_counts.get("text", 0) / 12.0, 4),
        round(type_counts.get("decoration", 0) / 8.0, 4),
        round(type_counts.get("drawing", 0) / 3.0, 4),
        round(len(flat) / 20.0, 4),
    ]


def _is_placeholder_shape(shape: object) -> bool:
    with contextlib.suppress(Exception):
        if bool(getattr(shape, "is_placeholder", False)):
            return True
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PLACEHOLDER


def _infer_element_type(
    *,
    shape: object,
    shape_type: object,
    has_text: bool,
    text: str,
    is_picture: bool,
    width: float,
    height: float,
    page_width: float,
    page_height: float,
) -> ReferenceElementType:
    if _is_placeholder_shape(shape) or shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return ReferenceElementType.PLACEHOLDER
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return ReferenceElementType.GROUP
    if shape_type == MSO_SHAPE_TYPE.TABLE or getattr(shape_type, "name", "") == "TABLE":
        return ReferenceElementType.TABLE
    if shape_type == MSO_SHAPE_TYPE.CHART:
        return ReferenceElementType.CHART
    if is_picture or shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
        # Drawing vs photo decided later by DrawingInferenceService.
        return ReferenceElementType.IMAGE
    if has_text and text.strip():
        return ReferenceElementType.TEXT
    page_area = max(page_width * page_height, 0.01)
    if width * height >= page_area * _DECORATION_AREA_RATIO:
        return ReferenceElementType.DECORATION
    if not has_text:
        return ReferenceElementType.DECORATION
    return ReferenceElementType.SHAPE


def _picture_alt_text(shape: object) -> str:
    """Read OOXML cNvPr descr/title — the usual home for picture alt text."""
    with contextlib.suppress(Exception):
        element = getattr(shape, "_element", None)
        if element is None:
            return ""
        for attr_path in ("nvPicPr", "nvSpPr"):
            container = getattr(element, attr_path, None)
            if container is None:
                continue
            c_nv_pr = getattr(container, "cNvPr", None)
            if c_nv_pr is None:
                continue
            descr = (c_nv_pr.get("descr") or "").strip()
            if descr:
                return descr
            title = (c_nv_pr.get("title") or "").strip()
            if title:
                return title
    return ""


_HARD_EDIT_SHAPE_NOTES: dict[object, str] = {
    MSO_SHAPE_TYPE.DIAGRAM: "hard_edit:smartart",
    MSO_SHAPE_TYPE.IGX_GRAPHIC: "hard_edit:smartart",
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "hard_edit:ole_embedded",
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "hard_edit:ole_linked",
    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT: "hard_edit:ole_control",
    MSO_SHAPE_TYPE.MEDIA: "hard_edit:media",
}


def _hard_edit_notes_for_shape(
    shape: object,
    *,
    shape_type: object,
    is_picture: bool,
    width: float,
    height: float,
    page_width: float,
    page_height: float,
) -> list[str]:
    """Tag shapes that are typically hard to replace in edit-based generation."""
    notes: list[str] = []
    hard = _HARD_EDIT_SHAPE_NOTES.get(shape_type)
    if hard:
        notes.append(hard)
    type_name = getattr(shape_type, "name", "") or ""
    if ("DIAGRAM" in type_name or "IGX" in type_name) and "hard_edit:smartart" not in notes:
        notes.append("hard_edit:smartart")
    if "OLE" in type_name and not any(n.startswith("hard_edit:ole") for n in notes):
        notes.append("hard_edit:ole_embedded")

    page_area = max(page_width * page_height, 0.01)
    if is_picture and width * height >= page_area * 0.85:
        notes.append("hard_edit:full_page_background")

    notes.extend(_shape_lock_notes(shape))
    if is_picture:
        notes.extend(_picture_complexity_notes(shape))
    return notes


def _shape_lock_notes(shape: object) -> list[str]:
    """Detect OOXML lock flags (noSelect / noGrp / noChangeAspect, …)."""
    notes: list[str] = []
    with contextlib.suppress(Exception):
        element = getattr(shape, "_element", None)
        if element is None:
            return notes
        for locks in element.iter():
            tag = getattr(locks, "tag", "") or ""
            if not str(tag).endswith("Locks"):
                continue
            attrs = getattr(locks, "attrib", {}) or {}
            restrictive = {
                "noSelect",
                "noGrp",
                "noMove",
                "noResize",
                "noRot",
                "noChangeAspect",
                "noEditPoints",
                "noAdjustHandles",
            }
            local_attrs = {str(k).split("}")[-1] for k in attrs}
            if local_attrs & restrictive:
                notes.append("hard_edit:locked")
                break
    return notes


def _picture_complexity_notes(shape: object) -> list[str]:
    """Crop / soft-edge / duotone style cues that complicate clean replacement."""
    notes: list[str] = []
    with contextlib.suppress(Exception):
        element = getattr(shape, "_element", None)
        if element is None:
            return notes
        xml = str(getattr(element, "xml", "") or "")
        lower = xml.lower()
        if "srcrect" in lower:
            notes.append("hard_edit:picture_crop")
        # blipFill alone is normal; only flag when effect-like tokens appear.
        if any(
            token in lower for token in ("alphamodfix", "duotone", "softedge", "blipfill")
        ) and any(
            token in lower
            for token in ("alphamodfix", "duotone", "softedge", "lumimod", "grayscl")
        ):
            notes.append("hard_edit:picture_effects")
    return notes


def _infer_semantic_role(
    *,
    element_type: ReferenceElementType,
    text: str,
    y: float,
    height: float,
    page_height: float,
    font_size_pt: float | None,
) -> str:
    lower = text.lower()
    if element_type == ReferenceElementType.PLACEHOLDER:
        return "placeholder"
    if element_type == ReferenceElementType.IMAGE:
        return "supporting_image" if height < page_height * 0.45 else "hero_image"
    if element_type == ReferenceElementType.DRAWING:
        return "drawing"
    if element_type == ReferenceElementType.CHART:
        return "chart"
    if element_type == ReferenceElementType.TABLE:
        return "table"
    if element_type == ReferenceElementType.GROUP:
        return "group"
    if element_type == ReferenceElementType.DECORATION:
        return "decoration"
    if y <= page_height * 0.22 and (font_size_pt or 0) >= 20:
        return "title"
    if y <= page_height * 0.35 and (font_size_pt or 0) >= 14:
        return "subtitle"
    if "来源" in text or "source" in lower:
        return "source"
    if any(token in text for token in ("% ", "㎡", "m²", "指标", "面积")):
        return "metric"
    if y >= page_height * 0.82:
        return "caption"
    return "body"


def _extension_for_image(content_type: str, blob: bytes) -> str:
    ctype = (content_type or "").lower()
    if "png" in ctype:
        return ".png"
    if "jpeg" in ctype or "jpg" in ctype:
        return ".jpg"
    if "gif" in ctype:
        return ".gif"
    if "webp" in ctype:
        return ".webp"
    if "emf" in ctype:
        return ".emf"
    if "wmf" in ctype:
        return ".wmf"
    if blob.startswith(b"\x89PNG"):
        return ".png"
    if blob[:2] == b"\xff\xd8":
        return ".jpg"
    return ".bin"


def _structural_signature(
    *,
    element_type: ReferenceElementType,
    x: float,
    y: float,
    width: float,
    height: float,
    font_name: str | None,
    fill_color: str | None,
    text: str,
    asset_hash: str,
    layout_name: str | None,
    master_name: str | None,
    placeholder: bool,
) -> str:
    text_hash = (
        hashlib.sha256(text.strip().encode("utf-8")).hexdigest()[:12] if text.strip() else ""
    )
    payload = {
        "type": element_type.value,
        "geom": (
            round(x, 2),
            round(y, 2),
            round(width, 2),
            round(height, 2),
        ),
        "font": font_name or "",
        "fill": fill_color or "",
        "text_hash": text_hash,
        "asset_hash": asset_hash or "",
        "layout": layout_name or "",
        "master": master_name or "",
        "placeholder": placeholder,
    }
    raw = json.dumps(payload, sort_keys=True, ensure_ascii=False)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:20]


def _mark_repeated_elements(slides: list[ReferenceSlideSnapshot]) -> None:
    """Mark cross-page repeats using structural signatures — not PowerPoint auto-names."""
    if len(slides) < 2:
        return
    threshold = max(2, len(slides) // 3)
    sig_pages: dict[str, set[int]] = defaultdict(set)
    for slide in slides:
        for element in slide.iter_elements():
            if not element.structural_signature:
                continue
            sig_pages[element.structural_signature].add(slide.slide_index)

    repeated = {sig for sig, pages in sig_pages.items() if len(pages) >= threshold}
    for slide in slides:
        for element in slide.iter_elements():
            if element.structural_signature not in repeated:
                continue
            element.repeats_across_pages = True
            # Only promote to decoration when the signature is chrome-like, not body copy.
            if _looks_like_chrome(element):
                element.likely_background_or_decoration = True
                if element.element_type == ReferenceElementType.TEXT and not element.text.strip():
                    element.element_type = ReferenceElementType.DECORATION
                    element.semantic_role = "decoration"


def _looks_like_chrome(element: ReferenceElement) -> bool:
    if element.element_type in {
        ReferenceElementType.DECORATION,
        ReferenceElementType.SHAPE,
        ReferenceElementType.PLACEHOLDER,
    }:
        return True
    if element.element_type == ReferenceElementType.TEXT:
        # Repeated empty/short labels near edges are often footer chrome.
        return len(element.text.strip()) <= 12 and element.semantic_role in {
            "",
            "body",
            "caption",
            "source",
        }
    if element.element_type in {
        ReferenceElementType.IMAGE,
        ReferenceElementType.DRAWING,
    }:
        # Shared identical image hash across pages can be logo/watermark.
        return bool(element.asset_id) and element.width * element.height < 4.0
    return False


class ReferencePptxParser:
    """Parse a reference PPTX into portable ReferencePresentation snapshots."""

    def __init__(self, *, drawing_inference: DrawingInferenceService | None = None) -> None:
        self._drawing_inference = drawing_inference or DrawingInferenceService()

    def parse(
        self,
        pptx_path: Path | str,
        *,
        workspace_dir: Path,
        name: str | None = None,
        capture_screenshots: bool = True,
    ) -> ReferencePresentation:
        source = Path(pptx_path)
        workspace = Path(workspace_dir)
        workspace.mkdir(parents=True, exist_ok=True)
        slides_dir = workspace / "slides"
        slides_dir.mkdir(parents=True, exist_ok=True)
        (workspace / "assets").mkdir(parents=True, exist_ok=True)

        warnings: list[str] = []
        try:
            presentation = Presentation(str(source))
        except Exception as exc:  # noqa: BLE001
            return ReferencePresentation(
                name=name or source.stem or "unreadable",
                source_filename=source.name,
                warnings=[f"无法打开 PPTX：{exc}"],
            )

        slide_width = int(presentation.slide_width or 0)
        slide_height = int(presentation.slide_height or 0)
        page_width = _emu_to_inches(slide_width) or 10.0
        page_height = _emu_to_inches(slide_height) or 5.625

        screenshot_map: dict[int, Path] = {}
        if capture_screenshots and screenshot_tools_available():
            pngs = export_pptx_slide_pngs(source, slides_dir)
            for index, png in enumerate(pngs):
                target = slides_dir / f"{_slide_id(index)}.png"
                if png.resolve() != target.resolve():
                    with contextlib.suppress(Exception):
                        if target.exists():
                            target.unlink()
                        png.replace(target)
                screenshot_map[index] = target if target.exists() else png
        elif capture_screenshots:
            warnings.append("截图工具不可用；仍输出结构 JSON，image_path 可为空。")

        slides: list[ReferenceSlideSnapshot] = []
        all_fonts: Counter[str] = Counter()
        all_colors: Counter[str] = Counter()

        for index, slide in enumerate(presentation.slides):
            try:
                snapshot = self._parse_slide(
                    slide,
                    index=index,
                    page_width=page_width,
                    page_height=page_height,
                    workspace=workspace,
                    image_path=_relative_slide_image(index)
                    if index in screenshot_map
                    else "",
                    fonts_counter=all_fonts,
                    colors_counter=all_colors,
                )
                slides.append(snapshot)
            except Exception as exc:  # noqa: BLE001 — page failure must not drop deck
                warnings.append(f"页面 {index + 1} 解析失败：{exc}")
                slides.append(
                    ReferenceSlideSnapshot(
                        slide_index=index,
                        slide_id=_slide_id(index),
                        width=page_width,
                        height=page_height,
                        parse_warnings=[str(exc)],
                        content_signature=f"parse_failed_{index}",
                    )
                )

        _mark_repeated_elements(slides)

        from archium.application.visual.induction_screenshot_embedding import (
            enrich_slide_screenshot_embeddings,
        )

        slides, attached = enrich_slide_screenshot_embeddings(slides, workspace, enabled=True)
        if attached:
            warnings.append(f"screenshot_embedding attached for {attached} slides")

        return ReferencePresentation(
            name=name or source.stem or "reference",
            source_filename=source.name,
            slide_count=len(slides),
            page_width=page_width,
            page_height=page_height,
            fonts=[n for n, _ in all_fonts.most_common(12)],
            colors=[c for c, _ in all_colors.most_common(12)],
            slides=slides,
            warnings=warnings,
            source_pptx_relative=_assert_relative("source.pptx"),
        )

    def _parse_slide(
        self,
        slide: object,
        *,
        index: int,
        page_width: float,
        page_height: float,
        workspace: Path,
        image_path: str,
        fonts_counter: Counter[str],
        colors_counter: Counter[str],
    ) -> ReferenceSlideSnapshot:
        layout_name: str | None = None
        master_name: str | None = None
        with contextlib.suppress(Exception):
            layout = getattr(slide, "slide_layout", None)
            if layout is not None:
                layout_name = getattr(layout, "name", None)
                master = getattr(layout, "slide_master", None)
                if master is not None:
                    master_name = getattr(master, "name", None)

        elements: list[ReferenceElement] = []
        text_content: list[str] = []
        image_assets: list[ReferenceAsset] = []
        page_fonts: list[str] = []
        page_colors: list[str] = []
        parse_warnings: list[str] = []
        image_seq = {"n": 0}

        shapes = list(getattr(slide, "shapes", []))
        for z_index, shape in enumerate(shapes):
            try:
                element, assets, texts, fonts, colors = self._parse_shape(
                    shape,
                    slide_index=index,
                    z_index=z_index,
                    id_prefix=f"s{index + 1:03d}_e{z_index + 1:03d}",
                    page_width=page_width,
                    page_height=page_height,
                    workspace=workspace,
                    layout_name=layout_name,
                    master_name=master_name,
                    image_seq=image_seq,
                )
            except Exception as exc:  # noqa: BLE001
                parse_warnings.append(f"shape[{z_index}] failed: {exc}")
                continue
            if element is None:
                continue
            elements.append(element)
            text_content.extend(texts)
            page_fonts.extend(fonts)
            page_colors.extend(colors)
            for font in fonts:
                fonts_counter[font] += 1
            for color in colors:
                colors_counter[color] += 1
            image_assets.extend(assets)

        notes = ""
        with contextlib.suppress(Exception):
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide is not None and notes_slide.notes_text_frame is not None:
                notes = (notes_slide.notes_text_frame.text or "").strip()

        # Drawing inference on flattened picture nodes (including group children).
        flat_for_inference = [n for e in elements for n in e.iter_self_and_descendants()]
        self._drawing_inference.refine_slide_elements(flat_for_inference, slide_notes=notes)

        image_count = sum(
            1
            for e in flat_for_inference
            if e.element_type
            in {ReferenceElementType.IMAGE, ReferenceElementType.DRAWING}
        )
        chart_count = sum(
            1 for e in flat_for_inference if e.element_type == ReferenceElementType.CHART
        )
        table_count = sum(
            1 for e in flat_for_inference if e.element_type == ReferenceElementType.TABLE
        )
        text_length = sum(len(t) for t in text_content)
        type_list = sorted(e.element_type.value for e in flat_for_inference)
        signature = _content_signature(
            element_types=type_list,
            image_count=image_count,
            text_length=text_length,
            chart_count=chart_count,
            table_count=table_count,
        )
        embedding = _visual_embedding(
            width=page_width,
            height=page_height,
            elements=elements,
            image_count=image_count,
            text_length=text_length,
            chart_count=chart_count,
            table_count=table_count,
        )

        return ReferenceSlideSnapshot(
            slide_index=index,
            slide_id=_slide_id(index),
            image_path=_assert_relative(image_path),
            width=page_width,
            height=page_height,
            master_name=master_name,
            layout_name=layout_name,
            elements=elements,
            text_content=text_content,
            image_assets=image_assets,
            visual_embedding=embedding,
            content_signature=signature,
            notes=notes,
            fonts=list(dict.fromkeys(page_fonts)),
            colors=list(dict.fromkeys(page_colors)),
            parse_warnings=parse_warnings,
        )

    def _parse_shape(
        self,
        shape: object,
        *,
        slide_index: int,
        z_index: int,
        id_prefix: str,
        page_width: float,
        page_height: float,
        workspace: Path,
        layout_name: str | None,
        master_name: str | None,
        image_seq: dict[str, int],
        depth: int = 0,
    ) -> tuple[
        ReferenceElement | None,
        list[ReferenceAsset],
        list[str],
        list[str],
        list[str],
    ]:
        try:
            x = _emu_to_inches(getattr(shape, "left", 0))
            y = _emu_to_inches(getattr(shape, "top", 0))
            width = max(_emu_to_inches(getattr(shape, "width", 0)), 0.05)
            height = max(_emu_to_inches(getattr(shape, "height", 0)), 0.05)
        except Exception:
            return None, [], [], [], []

        # Slides may bleed slightly outside the artboard; clamp for persistence.
        x = max(0.0, x)
        y = max(0.0, y)

        shape_type = getattr(shape, "shape_type", None)
        is_picture = shape_type in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        }
        # Placeholder picture hosts still expose .image
        if not is_picture:
            with contextlib.suppress(Exception):
                if _is_placeholder_shape(shape) and getattr(shape, "image", None) is not None:
                    is_picture = True

        has_text = bool(getattr(shape, "has_text_frame", False))
        text = ""
        font_name: str | None = None
        font_size_pt: float | None = None
        fill_color: str | None = None
        fonts: list[str] = []
        colors: list[str] = []
        texts: list[str] = []

        if has_text:
            lines: list[str] = []
            for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                chunk = (paragraph.text or "").strip()
                if chunk:
                    lines.append(chunk)
                for run in paragraph.runs:
                    if run.font.name:
                        font_name = run.font.name
                        fonts.append(run.font.name)
                    if run.font.size is not None:
                        with contextlib.suppress(Exception):
                            font_size_pt = float(run.font.size.pt)
                    try:
                        if run.font.color is not None and run.font.color.rgb is not None:
                            hex_color = _rgb_to_hex(run.font.color.rgb)
                            if hex_color:
                                colors.append(hex_color)
                    except Exception:
                        pass
            text = "\n".join(lines)
            if text:
                texts.append(text)

        with contextlib.suppress(Exception):
            fill = getattr(shape, "fill", None)
            if fill is not None and getattr(fill, "type", None) is not None:
                fore = getattr(fill, "fore_color", None)
                rgb = getattr(fore, "rgb", None) if fore is not None else None
                if rgb is not None:
                    fill_color = _rgb_to_hex(rgb)
                    if fill_color:
                        colors.append(fill_color)

        is_group = shape_type == MSO_SHAPE_TYPE.GROUP
        # Skip tiny marks (but keep groups / pictures / placeholders).
        if (
            not is_group
            and not is_picture
            and not _is_placeholder_shape(shape)
            and (width < 0.12 or height < 0.1)
        ):
            return None, [], texts, fonts, colors
        if (
            not is_group
            and not is_picture
            and not _is_placeholder_shape(shape)
            and not has_text
            and width * height < 0.25
        ):
            return None, [], texts, fonts, colors

        element_type = _infer_element_type(
            shape=shape,
            shape_type=shape_type,
            has_text=bool(text),
            text=text,
            is_picture=is_picture,
            width=width,
            height=height,
            page_width=page_width,
            page_height=page_height,
        )
        role = _infer_semantic_role(
            element_type=element_type,
            text=text,
            y=y,
            height=height,
            page_height=page_height,
            font_size_pt=font_size_pt,
        )
        shape_name = getattr(shape, "name", "") or ""
        alt_text = _picture_alt_text(shape) if (
            is_picture or element_type == ReferenceElementType.PLACEHOLDER
        ) else ""
        style_notes: list[str] = []
        placeholder_idx: int | None = None
        placeholder_type_raw = ""
        if element_type == ReferenceElementType.PLACEHOLDER:
            with contextlib.suppress(Exception):
                fmt = getattr(shape, "placeholder_format", None)
                if fmt is not None:
                    placeholder_type_raw = getattr(fmt, "type", "") or ""
                    style_notes.append(f"placeholder_type:{placeholder_type_raw}")
                    raw_idx = getattr(fmt, "idx", None)
                    if raw_idx is not None:
                        with contextlib.suppress(TypeError, ValueError):
                            placeholder_idx = int(raw_idx)
                        style_notes.append(f"placeholder_idx:{raw_idx}")
            if is_picture:
                style_notes.append("placeholder_hosts_picture")
            elif text:
                style_notes.append("placeholder_hosts_text")
        if is_group:
            style_notes.append("hard_edit:group")
        style_notes.extend(
            _hard_edit_notes_for_shape(
                shape,
                shape_type=shape_type,
                is_picture=is_picture,
                width=width,
                height=height,
                page_width=page_width,
                page_height=page_height,
            )
        )
        # De-dupe while preserving order.
        style_notes = list(dict.fromkeys(style_notes))

        assets: list[ReferenceAsset] = []
        asset_id: str | None = None
        asset_hash = ""
        if is_picture:
            image_seq["n"] += 1
            asset_id = f"refimg_{slide_index + 1:03d}_{image_seq['n']:03d}"
            asset, warning = self._extract_picture_asset(
                shape,
                slide_index=slide_index,
                image_index=image_seq["n"],
                asset_id=asset_id,
                width=width,
                height=height,
                workspace=workspace,
            )
            if warning:
                style_notes.append(warning)
            if asset is not None:
                assets.append(asset)
                asset_hash = asset.content_hash
                if alt_text:
                    style_notes.append(f"alt:{alt_text}")

        children: list[ReferenceElement] = []
        if is_group and depth < 6:
            child_shapes = list(getattr(shape, "shapes", []))
            for child_index, child in enumerate(child_shapes):
                child_element, child_assets, child_texts, child_fonts, child_colors = (
                    self._parse_shape(
                        child,
                        slide_index=slide_index,
                        z_index=child_index,
                        id_prefix=f"{id_prefix}_c{child_index + 1:03d}",
                        page_width=page_width,
                        page_height=page_height,
                        workspace=workspace,
                        layout_name=layout_name,
                        master_name=master_name,
                        image_seq=image_seq,
                        depth=depth + 1,
                    )
                )
                if child_element is not None:
                    children.append(child_element)
                assets.extend(child_assets)
                texts.extend(child_texts)
                fonts.extend(child_fonts)
                colors.extend(child_colors)

        signature = _structural_signature(
            element_type=element_type,
            x=x,
            y=y,
            width=width,
            height=height,
            font_name=font_name,
            fill_color=fill_color,
            text=text,
            asset_hash=asset_hash,
            layout_name=layout_name,
            master_name=master_name,
            placeholder=element_type == ReferenceElementType.PLACEHOLDER,
        )
        if shape_name and not _GENERIC_SHAPE_NAME.match(shape_name.strip()):
            style_notes.append(f"named:{shape_name}")

        placeholder_binding = None
        if element_type == ReferenceElementType.PLACEHOLDER:
            placeholder_binding = build_placeholder_binding_signature(
                placeholder_idx=placeholder_idx,
                placeholder_name=shape_name,
                placeholder_type=placeholder_type_raw,
                semantic_role=role,
                x=x,
                y=y,
                width=width,
                height=height,
                hosts_picture=is_picture,
                page_height=page_height,
            )
            # Prefer native placeholder type over the generic "placeholder" role.
            if placeholder_binding.semantic_role:
                role = placeholder_binding.semantic_role

        element = ReferenceElement(
            id=id_prefix,
            element_type=element_type,
            x=x,
            y=y,
            width=width,
            height=height,
            z_index=z_index,
            text=text,
            font_name=font_name,
            font_size_pt=font_size_pt,
            fill_color=fill_color,
            style_notes=style_notes,
            semantic_role=role,
            likely_background_or_decoration=element_type == ReferenceElementType.DECORATION,
            asset_id=asset_id,
            source_shape_name=shape_name,
            alt_text=alt_text,
            structural_signature=signature,
            placeholder_binding=placeholder_binding,
            children=children,
            parse_ok=True,
        )
        return element, assets, texts, fonts, colors

    def _extract_picture_asset(
        self,
        shape: object,
        *,
        slide_index: int,
        image_index: int,
        asset_id: str,
        width: float,
        height: float,
        workspace: Path,
    ) -> tuple[ReferenceAsset | None, str]:
        blob = b""
        content_type = ""
        with contextlib.suppress(Exception):
            image = shape.image  # type: ignore[attr-defined]
            blob = getattr(image, "blob", b"") or b""
            content_type = getattr(image, "content_type", "") or ""
        if not blob:
            return (
                ReferenceAsset(
                    id=asset_id,
                    asset_origin=REFERENCE_TEMPLATE_ASSET_ORIGIN,
                    relative_path="",
                    width=width,
                    height=height,
                    content_hash="",
                    content_type=content_type,
                    notes="reference_template_missing_blob",
                ),
                "picture_blob_missing",
            )

        content_hash = hashlib.sha256(blob).hexdigest()[:16]
        ext = _extension_for_image(content_type, blob)
        relative = (
            f"assets/{_slide_id(slide_index)}/image_{image_index:03d}{ext}"
        )
        relative = _assert_relative(relative)
        dest = workspace / relative
        dest.parent.mkdir(parents=True, exist_ok=True)
        dest.write_bytes(blob)
        return (
            ReferenceAsset(
                id=asset_id,
                asset_origin=REFERENCE_TEMPLATE_ASSET_ORIGIN,
                relative_path=relative,
                width=width,
                height=height,
                content_hash=content_hash,
                content_type=content_type or f"image/{ext.lstrip('.')}",
                notes="reference_template_only",
            ),
            "",
        )


# Keep Emu imported for parity with structure extractor / type checkers.
_ = Emu
