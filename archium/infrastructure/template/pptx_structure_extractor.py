"""Extract structure, fonts, and colors from a reference PPTX."""

from __future__ import annotations

import contextlib
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from archium.domain.visual.architectural_template import (
    ArchitecturalTemplateLayout,
    PowerPointMasterMetadata,
    TemplatePageType,
    TemplateSlot,
    TemplateSlotRole,
)

_EMU_PER_INCH = 914400.0


@dataclass
class ExtractedPageStructure:
    page_index: int
    page_width: float
    page_height: float
    slots: list[TemplateSlot] = field(default_factory=list)
    fonts: list[str] = field(default_factory=list)
    shape_colors: list[str] = field(default_factory=list)
    text_snippets: list[str] = field(default_factory=list)
    image_count: int = 0
    text_shape_count: int = 0


@dataclass
class PptxStructureExtraction:
    metadata: PowerPointMasterMetadata
    pages: list[ExtractedPageStructure]
    fonts: list[str]
    colors: list[str]
    warnings: list[str] = field(default_factory=list)


def _emu_to_inches(value: int | float | None) -> float:
    if value is None:
        return 0.0
    return round(float(value) / _EMU_PER_INCH, 4)


def _rgb_to_hex(rgb: object) -> str | None:
    try:
        return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"  # type: ignore[index]
    except Exception:
        return None


def _infer_slot_role(
    *,
    has_text: bool,
    text: str,
    is_picture: bool,
    y: float,
    height: float,
    page_height: float,
    font_size_pt: float | None,
) -> TemplateSlotRole:
    if is_picture:
        if height >= page_height * 0.45:
            return TemplateSlotRole.HERO_IMAGE
        return TemplateSlotRole.SUPPORTING_IMAGE
    if not has_text:
        return TemplateSlotRole.DECORATION
    lower = text.lower()
    if y <= page_height * 0.22 and (font_size_pt or 0) >= 20:
        return TemplateSlotRole.TITLE
    if y <= page_height * 0.35 and (font_size_pt or 0) >= 14:
        return TemplateSlotRole.SUBTITLE
    if "来源" in text or "source" in lower:
        return TemplateSlotRole.SOURCE
    if any(token in text for token in ("% ", "㎡", "m²", "指标", "面积")):
        return TemplateSlotRole.METRIC
    if y >= page_height * 0.82:
        return TemplateSlotRole.CAPTION
    return TemplateSlotRole.BODY


def _classify_page(page: ExtractedPageStructure) -> tuple[TemplatePageType, float, str]:
    roles = {slot.role for slot in page.slots}
    text_blob = " ".join(page.text_snippets)
    if page.page_index == 0 and TemplateSlotRole.TITLE in roles:
        return TemplatePageType.COVER, 0.7, "first page with title"
    if page.image_count >= 4:
        return TemplatePageType.PHOTO_GRID, 0.75, "multiple images"
    if page.image_count >= 2 and any("前后" in t or "before" in t.lower() for t in page.text_snippets):
        return TemplatePageType.BEFORE_AFTER, 0.7, "before/after cues"
    if page.image_count >= 2 and TemplateSlotRole.TITLE in roles:
        return TemplatePageType.CASE_COMPARISON, 0.55, "multi-image comparison candidate"
    if TemplateSlotRole.DRAWING in roles or any("平面" in t or "总图" in t for t in page.text_snippets):
        return TemplatePageType.DRAWING_FOCUS, 0.65, "drawing cues"
    if TemplateSlotRole.METRIC in roles or sum(ch.isdigit() for ch in text_blob) > 20:
        return TemplatePageType.METRIC, 0.6, "metric cues"
    if page.text_shape_count >= 3 and page.image_count == 0:
        return TemplatePageType.TEXT_ARGUMENT, 0.65, "text-heavy page"
    if "目录" in text_blob or "agenda" in text_blob.lower():
        return TemplatePageType.AGENDA, 0.8, "agenda keyword"
    if page.page_index > 0 and len(page.text_snippets) <= 2 and page.image_count <= 1:
        return TemplatePageType.SECTION, 0.45, "sparse section candidate"
    return TemplatePageType.UNKNOWN, 0.2, "insufficient signals"


class PptxStructureExtractor:
    """Extract page geometry, fonts, and colors from a PPTX for Template Studio."""

    def extract(self, pptx_path: Path) -> PptxStructureExtraction:
        warnings: list[str] = []
        try:
            presentation = Presentation(str(pptx_path))
        except Exception as exc:  # noqa: BLE001 — surface as soft extraction failure
            return PptxStructureExtraction(
                metadata=PowerPointMasterMetadata(encrypted_or_unreadable=True, notes=str(exc)),
                pages=[],
                fonts=[],
                colors=[],
                warnings=[f"无法打开 PPTX：{exc}"],
            )

        slide_width = int(presentation.slide_width or 0)
        slide_height = int(presentation.slide_height or 0)
        page_width = _emu_to_inches(slide_width) or 10.0
        page_height = _emu_to_inches(slide_height) or 5.625

        master_count = 0
        layout_count = 0
        try:
            master_count = len(presentation.slide_masters)
            layout_count = sum(len(master.slide_layouts) for master in presentation.slide_masters)
        except Exception:
            warnings.append("无法读取 Slide Master / Layout 元数据")

        pages: list[ExtractedPageStructure] = []
        all_fonts: Counter[str] = Counter()
        all_colors: Counter[str] = Counter()

        for index, slide in enumerate(presentation.slides):
            page = ExtractedPageStructure(
                page_index=index,
                page_width=page_width,
                page_height=page_height,
            )
            slot_index = 0
            for shape in slide.shapes:
                try:
                    x = _emu_to_inches(shape.left)
                    y = _emu_to_inches(shape.top)
                    width = max(_emu_to_inches(shape.width), 0.05)
                    height = max(_emu_to_inches(shape.height), 0.05)
                except Exception:
                    continue

                is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                has_text = bool(getattr(shape, "has_text_frame", False))
                text = ""
                font_name: str | None = None
                font_size_pt: float | None = None
                if has_text:
                    lines: list[str] = []
                    for paragraph in shape.text_frame.paragraphs:
                        chunk = (paragraph.text or "").strip()
                        if chunk:
                            lines.append(chunk)
                        for run in paragraph.runs:
                            if run.font.name:
                                font_name = run.font.name
                                all_fonts[run.font.name] += 1
                            if run.font.size is not None:
                                with contextlib.suppress(Exception):
                                    font_size_pt = float(run.font.size.pt)
                            try:
                                if run.font.color is not None and run.font.color.rgb is not None:
                                    hex_color = _rgb_to_hex(run.font.color.rgb)
                                    if hex_color:
                                        all_colors[hex_color] += 1
                                        page.shape_colors.append(hex_color)
                            except Exception:
                                pass
                    text = "\n".join(lines)
                    if text:
                        page.text_snippets.append(text)
                        page.text_shape_count += 1
                if is_picture:
                    page.image_count += 1

                if font_name:
                    page.fonts.append(font_name)

                # Skip tiny decorative marks.
                if width < 0.15 or height < 0.12:
                    continue
                if not has_text and not is_picture and width * height < 0.4:
                    # Keep larger shapes as decoration candidates.
                    continue

                role = _infer_slot_role(
                    has_text=bool(text),
                    text=text,
                    is_picture=is_picture,
                    y=y,
                    height=height,
                    page_height=page_height,
                    font_size_pt=font_size_pt,
                )
                if role == TemplateSlotRole.HERO_IMAGE and "平面" in text:
                    role = TemplateSlotRole.DRAWING

                slot_index += 1
                page.slots.append(
                    TemplateSlot(
                        id=f"p{index}_s{slot_index}",
                        role=role,
                        required=role in {TemplateSlotRole.TITLE, TemplateSlotRole.HERO_IMAGE},
                        x=x,
                        y=y,
                        width=width,
                        height=height,
                        accepted_node_types=["image", "drawing"] if is_picture else ["text"],
                        label=(text[:40] if text else role.value),
                        source_shape_name=getattr(shape, "name", "") or "",
                        auto_detected=True,
                    )
                )
            pages.append(page)

        fonts = [name for name, _ in all_fonts.most_common(12)]
        colors = [color for color, _ in all_colors.most_common(12)]
        metadata = PowerPointMasterMetadata(
            slide_count=len(pages),
            slide_width_emu=slide_width,
            slide_height_emu=slide_height,
            has_slide_master=master_count > 0,
            master_count=master_count,
            layout_count=layout_count,
        )
        return PptxStructureExtraction(
            metadata=metadata,
            pages=pages,
            fonts=fonts,
            colors=colors,
            warnings=warnings,
        )

    def to_layouts(
        self,
        extraction: PptxStructureExtraction,
        *,
        screenshot_paths: dict[int, Path] | None = None,
    ) -> list[ArchitecturalTemplateLayout]:
        layouts: list[ArchitecturalTemplateLayout] = []
        screenshots = screenshot_paths or {}
        for page in extraction.pages:
            page_type, confidence, notes = _classify_page(page)
            supports_photo = page.image_count > 0
            supports_drawing = any(slot.role == TemplateSlotRole.DRAWING for slot in page.slots)
            supports_metrics = any(slot.role == TemplateSlotRole.METRIC for slot in page.slots)
            layouts.append(
                ArchitecturalTemplateLayout(
                    name=f"页面 {page.page_index + 1}",
                    description=notes,
                    page_index=page.page_index,
                    page_type=page_type,
                    suitable_slide_types=[page_type.value],
                    slots=list(page.slots),
                    supports_drawing=supports_drawing,
                    supports_photo=supports_photo,
                    supports_metrics=supports_metrics,
                    supports_case_reference=page_type
                    in {TemplatePageType.CASE_COMPARISON, TemplatePageType.BEFORE_AFTER},
                    minimum_asset_count=1 if supports_photo or supports_drawing else 0,
                    maximum_asset_count=max(page.image_count, 1),
                    preview_image_path=str(screenshots.get(page.page_index, "")),
                    page_width=page.page_width,
                    page_height=page.page_height,
                    extracted_fonts=list(dict.fromkeys(page.fonts)),
                    extracted_colors=list(dict.fromkeys(page.shape_colors)),
                    classification_confidence=confidence,
                    classification_notes=notes,
                )
            )
        return layouts


# Keep Emu import referenced for type checkers / future master parsing.
_ = Emu
