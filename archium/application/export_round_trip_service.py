"""Export round-trip validation service — PPTX vs RenderScene QA."""

from __future__ import annotations

import hashlib
import tempfile
from pathlib import Path
from uuid import UUID

from sqlalchemy.orm import Session

from archium.application.visual.studio_scene_service import StudioSceneService
from archium.config.settings import Settings, get_settings
from archium.domain.export_round_trip import (
    ExportRoundTripReport,
    RoundTripStatus,
    SlideRoundTripResult,
)
from archium.domain.visual.render_scene import (
    DrawingNode,
    ImageNode,
    RenderScene,
    TextNode,
    compute_scene_hash,
)
from archium.infrastructure.renderers.canvas_renderer import CanvasRenderer
from archium.infrastructure.renderers.pptx_screenshot import (
    export_pptx_slide_pngs,
    screenshot_tools_available,
)
from archium.infrastructure.renderers.renderer_conformance import snapshot_from_scene
from archium.infrastructure.vision.screenshot_qa import (
    compare_png_pptx_screenshots,
    load_image,
)

_TEXT_RECALL_BLOCK_THRESHOLD = 0.5
_TEXT_RECALL_REVIEW_THRESHOLD = 0.95
_SIMILARITY_REVIEW_THRESHOLD = 0.85
_GEOMETRY_REVIEW_THRESHOLD = 0.8


class ExportRoundTripService:
    """Validate exported PPTX against source RenderScenes."""

    def __init__(
        self,
        session: Session,
        *,
        settings: Settings | None = None,
    ) -> None:
        self._session = session
        self._settings = settings or get_settings()
        self._scene_service = StudioSceneService(session, settings=self._settings)
        self._canvas = CanvasRenderer()

    def validate_pptx_export(
        self,
        *,
        presentation_id: UUID,
        pptx_path: Path | str,
        revision_id: UUID | None = None,
        export_file_hash: str | None = None,
    ) -> ExportRoundTripReport:
        pptx = Path(pptx_path)
        scene_results = self._scene_service.ensure_scenes_for_presentation(
            presentation_id,
            force_recompile=False,
        )
        scenes = [result.scene for result in scene_results]
        scenes_by_slide = {scene.slide_id: scene for scene in scenes}

        source_hash = _combined_scene_hash(scenes)
        file_hash = export_file_hash or _file_hash(pptx)

        slide_results: list[SlideRoundTripResult] = []
        all_missing_text: list[str] = []
        all_missing_assets: list[str] = []
        all_font_subs: list[str] = []
        all_drawing_issues: list[str] = []
        all_citation_issues: list[str] = []
        all_warnings: list[str] = []
        all_blockers: list[str] = []
        changed_origins: list[str] = []

        tools_ok = screenshot_tools_available()
        pptx_pngs: list[Path] = []
        preview_hash_parts: list[str] = []

        with tempfile.TemporaryDirectory(prefix="archium-roundtrip-") as tmp:
            tmp_dir = Path(tmp)
            if tools_ok and pptx.is_file():
                pptx_pngs = export_pptx_slide_pngs(pptx, tmp_dir)
            per_slide_pptx = _snapshots_by_slide(pptx) if pptx.is_file() else []

            if pptx.is_file():
                from pptx import Presentation

                presentation = Presentation(str(pptx))
                if len(presentation.slides) != len(scenes):
                    all_blockers.append(
                        f"页面数量不一致：源 Scene {len(scenes)} 页，导出 PPTX {len(presentation.slides)} 页"
                    )

            for index, scene in enumerate(scenes):
                pptx_snap = per_slide_pptx[index] if index < len(per_slide_pptx) else None
                source_snap = snapshot_from_scene(scene)
                slide_result = self._validate_slide(
                    scene,
                    slide_order=index,
                    pptx_snap=pptx_snap,
                    pptx_png=pptx_pngs[index] if index < len(pptx_pngs) else None,
                    preview_dir=tmp_dir,
                )
                slide_results.append(slide_result)
                all_missing_text.extend(slide_result.missing_text_nodes)
                all_missing_assets.extend(slide_result.missing_assets)
                all_font_subs.extend(slide_result.font_substitutions)
                all_drawing_issues.extend(slide_result.drawing_integrity_issues)
                all_citation_issues.extend(slide_result.citation_integrity_issues)
                all_warnings.extend(slide_result.warnings)

                for node in scene.nodes:
                    if isinstance(node, ImageNode) and node.asset_origin not in {
                        "project_upload",
                        "",
                    }:
                        marker = f"{node.id}:{node.asset_origin}"
                        if marker not in changed_origins:
                            changed_origins.append(marker)

            if pptx_pngs:
                for png_path in pptx_pngs:
                    preview_hash_parts.append(_file_hash(png_path))

        text_rates = [slide.text_match_rate for slide in slide_results] or [1.0]
        geom_rates = [slide.geometry_match_rate for slide in slide_results] or [1.0]
        sim_scores = [slide.similarity_score for slide in slide_results if slide.similarity_score >= 0]

        avg_text = sum(text_rates) / len(text_rates)
        avg_geom = sum(geom_rates) / len(geom_rates)
        avg_sim = sum(sim_scores) / len(sim_scores) if sim_scores else -1.0

        if avg_text < _TEXT_RECALL_BLOCK_THRESHOLD:
            all_blockers.append(
                f"文本召回率 {avg_text:.0%} 低于阻塞阈值 {_TEXT_RECALL_BLOCK_THRESHOLD:.0%}"
            )
        if all_drawing_issues:
            for issue in all_drawing_issues:
                if "fit_mode" in issue or "cover" in issue:
                    all_blockers.append(f"图纸保护：{issue}")

        status = _derive_status(
            text_match_rate=avg_text,
            geometry_match_rate=avg_geom,
            similarity_score=avg_sim,
            drawing_issues=all_drawing_issues,
            blockers=all_blockers,
            warnings=all_warnings,
        )

        if not tools_ok:
            all_warnings.append("未检测到 PPTX 截图工具，视觉相似度未计算。")
            if status == RoundTripStatus.PASS:
                status = RoundTripStatus.PASS_WITH_WARNINGS

        return ExportRoundTripReport(
            presentation_id=presentation_id,
            revision_id=revision_id,
            source_scene_hash=source_hash,
            export_file_hash=file_hash,
            rendered_preview_hash=hashlib.sha256(
                "|".join(preview_hash_parts).encode("utf-8")
            ).hexdigest()[:16]
            if preview_hash_parts
            else "",
            similarity_score=avg_sim,
            text_match_rate=avg_text,
            geometry_match_rate=avg_geom,
            missing_text_nodes=all_missing_text[:50],
            missing_assets=all_missing_assets[:50],
            changed_asset_origins=changed_origins[:30],
            font_substitutions=sorted(set(all_font_subs)),
            drawing_integrity_issues=all_drawing_issues[:50],
            citation_integrity_issues=all_citation_issues[:50],
            slides=slide_results,
            warnings=all_warnings,
            blockers=all_blockers,
            status=status,
            screenshot_tools_available=tools_ok,
        )

    def _validate_slide(
        self,
        scene: RenderScene,
        *,
        slide_order: int,
        pptx_snap,
        pptx_png: Path | None,
        preview_dir: Path,
    ) -> SlideRoundTripResult:
        source_snap = snapshot_from_scene(scene)
        missing_text: list[str] = []
        missing_assets: list[str] = []
        font_subs: list[str] = []
        drawing_issues: list[str] = []
        citation_issues: list[str] = []
        warnings: list[str] = []

        text_rate = 1.0
        geom_rate = 1.0
        similarity = -1.0

        if pptx_snap is not None:
            text_rate, missing_text = _text_recall(source_snap.text_values, pptx_snap.text_values)
            geom_rate = _geometry_match(source_snap, pptx_snap)

        for node in scene.nodes:
            if isinstance(node, (ImageNode, DrawingNode)) and node.asset_unresolved:
                missing_assets.append(node.id)

        for font in scene.font_assets:
            if font.resolved_family and font.resolved_family != font.family:
                font_subs.append(f"{font.family}→{font.resolved_family}")

        drawing_issues.extend(_drawing_integrity_checks(scene))
        citation_issues.extend(_citation_integrity_checks(scene))

        if pptx_png is not None and pptx_png.is_file():
            source_preview = preview_dir / f"source_{scene.slide_id}.png"
            self._canvas.render_preview(scene, source_preview)
            source_image = load_image(source_preview)
            pptx_image = load_image(pptx_png)
            if source_image is not None and pptx_image is not None:
                check = compare_png_pptx_screenshots(source_image, pptx_image)
                mse = float(check.evidence.get("mse", 9999))
                similarity = max(0.0, min(1.0, 1.0 - mse / 5000.0))
                if not check.passed:
                    warnings.append(check.description)

        return SlideRoundTripResult(
            slide_id=scene.slide_id,
            slide_order=slide_order,
            text_match_rate=text_rate,
            geometry_match_rate=geom_rate,
            similarity_score=similarity,
            missing_text_nodes=missing_text,
            missing_assets=missing_assets,
            font_substitutions=font_subs,
            drawing_integrity_issues=drawing_issues,
            citation_integrity_issues=citation_issues,
            warnings=warnings,
        )


def _combined_scene_hash(scenes: list[RenderScene]) -> str:
    if not scenes:
        return ""
    digest = hashlib.sha256()
    for scene in sorted(scenes, key=lambda item: str(item.slide_id)):
        digest.update(compute_scene_hash(scene).encode("utf-8"))
    return digest.hexdigest()[:16]


def _file_hash(path: Path) -> str:
    if not path.is_file():
        return ""
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(65536), b""):
            digest.update(chunk)
    return digest.hexdigest()[:16]


def _snapshots_by_slide(pptx_path: Path) -> list:
    from pptx import Presentation

    from archium.infrastructure.renderers.renderer_conformance import RendererSnapshot

    presentation = Presentation(str(pptx_path))
    snapshots: list[RendererSnapshot] = []
    for slide in presentation.slides:
        texts: list[str] = []
        image_count = 0
        for shape in slide.shapes:
            if getattr(shape, "text", "") and str(shape.text).strip():
                texts.append(str(shape.text).strip())
            if shape.shape_type == 13:
                image_count += 1
        snapshots.append(
            RendererSnapshot(
                text_values=tuple(texts),
                image_node_ids=tuple(str(i) for i in range(image_count)),
                node_count=len(texts) + image_count,
                background_color="",
            )
        )
    return snapshots


def _text_recall(
    source_texts: tuple[str, ...],
    exported_texts: tuple[str, ...],
) -> tuple[float, list[str]]:
    if not source_texts:
        return 1.0, []
    exported_blob = "\n".join(exported_texts)
    missing: list[str] = []
    matched = 0
    for text in source_texts:
        normalized = text.strip()
        if not normalized:
            continue
        if normalized in exported_blob or any(
            normalized in item or item in normalized for item in exported_texts
        ):
            matched += 1
        else:
            missing.append(normalized[:120])
    total = sum(1 for text in source_texts if text.strip())
    if total == 0:
        return 1.0, []
    return matched / total, missing


def _geometry_match(source, exported) -> float:
    if source.node_count <= 0:
        return 1.0
    ratio = exported.node_count / source.node_count
    if ratio > 1.0:
        ratio = 1.0 / ratio
    image_ref = len(source.image_node_ids)
    image_exp = len(exported.image_node_ids)
    if image_ref > 0:
        image_ratio = min(image_exp, image_ref) / image_ref
        return (ratio + image_ratio) / 2
    return ratio


def _drawing_integrity_checks(scene: RenderScene) -> list[str]:
    issues: list[str] = []
    for node in scene.nodes:
        if not isinstance(node, DrawingNode):
            continue
        if node.fit_mode not in {"contain", "safe_crop"}:
            issues.append(f"{node.id}:fit_mode={node.fit_mode}")
        if node.drawing_type in {"site_plan", "floor_plan", "circulation_plan"}:
            if not node.north_arrow_visible and not node.scale_label:
                issues.append(f"{node.id}:missing_north_or_scale")
        if node.crop_allowed and node.fit_mode == "contain":
            issues.append(f"{node.id}:crop_allowed_with_contain")
    return issues


def _citation_integrity_checks(scene: RenderScene) -> list[str]:
    issues: list[str] = []
    for node in scene.nodes:
        if isinstance(node, ImageNode) and node.asset_origin == "reference_case":
            if not node.caption_node_id:
                issues.append(f"{node.id}:reference_without_caption")
    return issues


def _derive_status(
    *,
    text_match_rate: float,
    geometry_match_rate: float,
    similarity_score: float,
    drawing_issues: list[str],
    blockers: list[str],
    warnings: list[str],
) -> RoundTripStatus:
    if blockers:
        return RoundTripStatus.BLOCKED
    if text_match_rate < _TEXT_RECALL_REVIEW_THRESHOLD:
        return RoundTripStatus.NEEDS_REVIEW
    if geometry_match_rate < _GEOMETRY_REVIEW_THRESHOLD:
        return RoundTripStatus.NEEDS_REVIEW
    if similarity_score >= 0 and similarity_score < _SIMILARITY_REVIEW_THRESHOLD:
        return RoundTripStatus.NEEDS_REVIEW
    if drawing_issues or warnings:
        return RoundTripStatus.PASS_WITH_WARNINGS
    return RoundTripStatus.PASS
