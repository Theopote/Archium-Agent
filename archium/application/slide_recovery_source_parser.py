"""Parse external page sources into proxy RenderScenes for slide recovery."""

from __future__ import annotations

import hashlib
import tempfile
from dataclasses import dataclass
from pathlib import Path
from uuid import uuid4

from archium.domain.visual.render_scene import (
    BackgroundStyle,
    ImageNode,
    RenderNode,
    RenderScene,
    ShapeNode,
    TextNode,
)
from archium.exceptions import WorkflowError
from archium.infrastructure.slide_recovery.pdf_page_renderer import render_pdf_page_png
from archium.infrastructure.slide_recovery.pptx_slide_renderer import render_pptx_slide_png

_EMU_PER_INCH = 914400.0
_SUPPORTED_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}


@dataclass(frozen=True)
class ParsedSourcePage:
    """One external page loaded for slide recovery."""

    scene: RenderScene
    page_id: str
    preview_image_path: Path | None = None
    source_kind: str = "image"


def parse_source_page(
    path: Path | str,
    *,
    slide_index: int = 0,
    workspace_dir: Path | None = None,
) -> ParsedSourcePage:
    """Load a single page from PNG/JPG, PDF, or PPTX into a proxy RenderScene."""
    source = Path(path)
    if not source.is_file():
        raise WorkflowError(f"源文件不存在：{source}")

    suffix = source.suffix.lower()
    if suffix in _SUPPORTED_IMAGE_SUFFIXES:
        scene, image_path = _parse_image_page(source, workspace_dir=workspace_dir)
        return ParsedSourcePage(
            scene=scene,
            page_id=source.stem,
            preview_image_path=image_path,
            source_kind="image",
        )
    if suffix == ".pdf":
        return _parse_pdf_page(
            source,
            page_index=slide_index,
            workspace_dir=workspace_dir,
        )
    if suffix == ".pptx":
        return _parse_pptx_slide(
            source,
            slide_index=slide_index,
            workspace_dir=workspace_dir,
        )
    raise WorkflowError(f"不支持的文件类型：{suffix}")


def _parse_pdf_page(
    path: Path,
    *,
    page_index: int,
    workspace_dir: Path | None,
) -> ParsedSourcePage:
    if workspace_dir is None:
        workspace_dir = Path(tempfile.gettempdir()) / "archium-slide-recovery" / "pdf"
    png_path, page_w, page_h = render_pdf_page_png(
        path,
        page_index=page_index,
        workspace_dir=workspace_dir,
    )
    storage_uri = _materialize_asset(png_path, workspace_dir=workspace_dir, label="pdf_page")
    scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=page_w,
        page_height=page_h,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=[
            ImageNode(
                id="source_page_image",
                x=0,
                y=0,
                width=page_w,
                height=page_h,
                z_index=0,
                storage_uri=storage_uri,
                semantic_role="source_page",
                asset_origin="project_upload",
                fit_mode="contain",
            ),
        ],
    )
    return ParsedSourcePage(
        scene=scene,
        page_id=f"{path.stem}_page_{page_index + 1:03d}",
        preview_image_path=png_path,
        source_kind="pdf",
    )


def _parse_image_page(
    path: Path,
    *,
    workspace_dir: Path | None,
) -> tuple[RenderScene, Path]:
    page_w, page_h = _image_page_size(path)
    storage_uri, image_path = _materialize_asset_with_path(
        path,
        workspace_dir=workspace_dir,
        label="page_image",
    )
    scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=page_w,
        page_height=page_h,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=[
            ImageNode(
                id="source_page_image",
                x=0,
                y=0,
                width=page_w,
                height=page_h,
                z_index=0,
                storage_uri=storage_uri,
                semantic_role="source_page",
                asset_origin="project_upload",
                fit_mode="contain",
            ),
        ],
    )
    return scene, image_path


def _parse_pptx_slide(
    path: Path,
    *,
    slide_index: int,
    workspace_dir: Path | None,
) -> ParsedSourcePage:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:  # pragma: no cover
        raise WorkflowError("需要 python-pptx 才能解析 PPTX 页面。") from exc

    presentation = Presentation(str(path))
    if slide_index < 0 or slide_index >= len(presentation.slides):
        raise WorkflowError(
            f"PPTX 只有 {len(presentation.slides)} 页，无法读取第 {slide_index + 1} 页。"
        )

    page_w = float(presentation.slide_width or 0) / _EMU_PER_INCH
    page_h = float(presentation.slide_height or 0) / _EMU_PER_INCH
    if page_w <= 0 or page_h <= 0:
        raise WorkflowError("PPTX 页面尺寸无效。")
    slide = presentation.slides[slide_index]
    nodes: list[RenderNode] = []
    z_index = 0

    for shape_index, shape in enumerate(slide.shapes):
        left = float(shape.left) / _EMU_PER_INCH
        top = float(shape.top) / _EMU_PER_INCH
        width = max(float(shape.width) / _EMU_PER_INCH, 0.05)
        height = max(float(shape.height) / _EMU_PER_INCH, 0.05)

        text = str(getattr(shape, "text", "") or "").strip()
        if text:
            nodes.append(
                TextNode(
                    id=f"text_{shape_index}",
                    x=left,
                    y=top,
                    width=width,
                    height=height,
                    z_index=z_index,
                    text=text,
                    semantic_role=_infer_text_role(text, left, top, page_w, page_h),
                    font_family="Microsoft YaHei",
                    font_size=18,
                    color="#1A1A1A",
                    line_height=24,
                )
            )
            z_index += 1
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            storage_uri = _extract_picture_asset(
                shape,
                source_path=path,
                slide_index=slide_index,
                shape_index=shape_index,
                workspace_dir=workspace_dir,
            )
            nodes.append(
                ImageNode(
                    id=f"image_{shape_index}",
                    x=left,
                    y=top,
                    width=width,
                    height=height,
                    z_index=z_index,
                    storage_uri=storage_uri,
                    semantic_role="image",
                    asset_origin="project_upload",
                    fit_mode="contain",
                )
            )
            z_index += 1
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            nodes.append(
                ShapeNode(
                    id=f"line_{shape_index}",
                    x=left,
                    y=top,
                    width=max(width, 0.01),
                    height=max(height, 0.01),
                    z_index=z_index,
                    shape_kind="line",
                    stroke_color="#333333",
                    stroke_width=1,
                    semantic_role="line",
                )
            )
            z_index += 1

    if not nodes:
        nodes.append(
            ImageNode(
                id="pptx_fallback",
                x=0,
                y=0,
                width=page_w,
                height=page_h,
                z_index=0,
                storage_uri=f"file://{path.resolve().as_posix()}",
                semantic_role="source_page",
                asset_origin="project_upload",
                fit_mode="contain",
            )
        )

    preview_dir = (
        workspace_dir / "pptx_previews" if workspace_dir is not None else None
    )
    preview_image = render_pptx_slide_png(
        path,
        slide_index=slide_index,
        workspace_dir=preview_dir,
    )

    scene = RenderScene(
        slide_id=uuid4(),
        layout_plan_id=uuid4(),
        page_width=page_w,
        page_height=page_h,
        background=BackgroundStyle(color="#FFFFFF"),
        nodes=nodes,
    )
    return ParsedSourcePage(
        scene=scene,
        page_id=f"{path.stem}_slide_{slide_index + 1:03d}",
        preview_image_path=preview_image,
        source_kind="pptx",
    )


def _infer_text_role(text: str, x: float, y: float, page_w: float, page_h: float) -> str:
    if y <= page_h * 0.2 and len(text) <= 40:
        return "title"
    if y >= page_h * 0.85:
        return "caption"
    if "\t" in text or text.count(" ") >= 3 and len(text) < 30:
        return "table_cell"
    return "body"


def _image_page_size(path: Path) -> tuple[float, float]:
    try:
        from PIL import Image
    except ImportError:  # pragma: no cover
        return 10.0, 5.625
    with Image.open(path) as image:
        width, height = image.size
    if width <= 0 or height <= 0:
        return 10.0, 5.625
    aspect = height / width
    page_w = 10.0
    return page_w, round(page_w * aspect, 4)


def _materialize_asset(
    path: Path,
    *,
    workspace_dir: Path | None,
    label: str,
) -> str:
    uri, _ = _materialize_asset_with_path(path, workspace_dir=workspace_dir, label=label)
    return uri


def _materialize_asset_with_path(
    path: Path,
    *,
    workspace_dir: Path | None,
    label: str,
) -> tuple[str, Path]:
    digest = hashlib.sha256(path.read_bytes()).hexdigest()[:16]
    if workspace_dir is not None:
        workspace_dir.mkdir(parents=True, exist_ok=True)
        target = workspace_dir / f"{label}_{digest}{path.suffix.lower() or '.bin'}"
        if not target.is_file():
            target.write_bytes(path.read_bytes())
        resolved = target.resolve()
        return f"file://{resolved.as_posix()}", resolved
    resolved = path.resolve()
    return f"file://{resolved.as_posix()}", resolved


def _extract_picture_asset(
    shape: object,
    *,
    source_path: Path,
    slide_index: int,
    shape_index: int,
    workspace_dir: Path | None,
) -> str:
    image = getattr(shape, "image", None)
    blob = getattr(image, "blob", None) if image is not None else None
    if not blob:
        return f"file://{source_path.resolve().as_posix()}#slide={slide_index}&shape={shape_index}"

    ext = getattr(getattr(image, "ext", None), "lower", lambda: "png")()
    if not ext:
        ext = "png"
    if workspace_dir is None:
        workspace_dir = Path(tempfile.gettempdir()) / "archium-slide-recovery"
    workspace_dir.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha256(blob).hexdigest()[:16]
    target = workspace_dir / f"slide_{slide_index + 1:03d}_img_{shape_index}_{digest}.{ext}"
    if not target.is_file():
        target.write_bytes(blob)
    return f"file://{target.resolve().as_posix()}"
